#!/usr/bin/env python3
"""
ChIP-seq Peak Calling Report Generator

Generates:
  1. PPT report with QC metrics, plots, and per-module summaries.
  2. Excel workbook with ALL result data from every module.

Usage:
    python generate_report.py \
        --samples Pop5IP --samples Rpp14IP --samples Rpp21IP \
        --input-samples Pop5Input --input-samples Rpp14Input --input-samples Rpp21Input \
        --ip-input-pair Pop5IP:Pop5Input --ip-input-pair Rpp14IP:Rpp14Input --ip-input-pair Rpp21IP:Rpp21Input \
        --peaks-dir ... --annotation-dir ... --qc-dir ... --log-dir ... --markdup-dir ... \
        --trim-dir ... --metrics-dir ... --te-dir ... \
        --output report.pptx --excel-output report.xlsx
"""

import argparse
import csv
import os
import sys
import tempfile
from collections import Counter

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ============================================================
# Layout & Colors
# ============================================================
SLIDE_W = 10.0
SLIDE_H = 5.625
HEADER_H = 0.65
MARGIN_L = 0.45
CONTENT_W = SLIDE_W - 2 * MARGIN_L
CONTENT_TOP = HEADER_H + 0.18
DPI = 300

C_NAVY = RGBColor(0x18, 0x25, 0x43)
C_ACCENT = RGBColor(0x00, 0x94, 0xD8)
C_GREEN = RGBColor(0x00, 0xA8, 0x78)
C_RED = RGBColor(0xD6, 0x45, 0x45)
C_ORANGE = RGBColor(0xF0, 0x9A, 0x36)
C_TEXT = RGBColor(0x33, 0x33, 0x33)
C_MUTED = RGBColor(0x74, 0x74, 0x74)
C_BG = RGBColor(0xF6, 0xF8, 0xFB)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_ROW_ALT = RGBColor(0xF1, 0xF5, 0xFA)

ALIGN_THRESH = 80.0
FRIP_THRESH = 0.20


# ============================================================
# I18N
# ============================================================
I18N = {
    "zh": {
        "report_title": "ChIP-seq Peak Calling 分析报告",
        "pipeline_label": "分析流程", "date_label": "报告日期",
        "workflow_title": "分析流程概览",
        "trim_title": "TrimGalore 质控",
        "alignment_title": "Bowtie2 比对统计",
        "markdup_title": "GATK MarkDuplicates 统计",
        "peak_calling_title": "MACS3 Peak Calling 结果",
        "annotation_title": "Peak 注释 — 基因组区域分布",
        "top_genes_title": "Top 靶基因",
        "summary_title": "总结与 QC 评估",
        "tss_title": "TSS 距离分布",
        "te_enrichment_title": "TE 富集分析",
        "peak_centric_te_title": "Peak 视角 TE 分析",
        "sample": "样本", "type": "类型", "ip_type": "IP", "input_type": "Input",
        "read_pairs": "Read Pairs", "dup_pairs": "Dup Pairs", "dup_rate": "Dup Rate",
        "peak_count": "Peak 数", "frip_score": "FRiP",
        "conclusions": "结论", "recommendations": "建议",
    },
    "en": {
        "report_title": "ChIP-seq Peak Calling Report",
        "pipeline_label": "Pipeline", "date_label": "Date",
        "workflow_title": "Workflow Overview",
        "trim_title": "TrimGalore QC",
        "alignment_title": "Bowtie2 Alignment Statistics",
        "markdup_title": "GATK MarkDuplicates Statistics",
        "peak_calling_title": "MACS3 Peak Calling Results",
        "annotation_title": "Peak Annotation — Genomic Region Distribution",
        "top_genes_title": "Top Target Genes",
        "summary_title": "Summary & QC Assessment",
        "tss_title": "TSS Distance Distribution",
        "te_enrichment_title": "TE Enrichment Analysis",
        "peak_centric_te_title": "Peak-Centric TE Analysis",
        "sample": "Sample", "type": "Type", "ip_type": "IP", "input_type": "Input",
        "read_pairs": "Read Pairs", "dup_pairs": "Dup Pairs", "dup_rate": "Dup Rate",
        "peak_count": "Peak Count", "frip_score": "FRiP",
        "conclusions": "Conclusions", "recommendations": "Recommendations",
    },
}


def t(key, lang="zh"):
    return I18N.get(lang, I18N["zh"]).get(key, key)


# ============================================================
# PPT helpers
# ============================================================

def _header(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(HEADER_H))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_NAVY; bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(MARGIN_L), Inches(0.08), Inches(CONTENT_W), Inches(0.45))
    p = tx.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = C_WHITE


def _bg(slide):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C_BG


def _textbox(slide, left, top, width, height, text, font_size=11, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size)
    p.font.bold = bold; p.font.color.rgb = color; p.alignment = align
    return box


def _bullets(slide, left, top, width, height, items, font_size=10, color=C_TEXT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"; p.font.size = Pt(font_size); p.font.color.rgb = color; p.space_after = Pt(3)
    return box


def _table(slide, left, top, width, height, data, font_size=10):
    rows, cols = len(data), len(data[0]) if data else 1
    shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    tbl = shape.table
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c); cell.text = str(val); cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = C_ACCENT
            elif r % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = C_ROW_ALT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.bold = (r == 0); p.font.color.rgb = C_WHITE if r == 0 else C_TEXT
    return tbl


def _add_img(slide, path, left, top, max_w, max_h):
    if path and os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(max_w), Inches(max_h))


# ============================================================
# Data Loaders
# ============================================================

def load_trim_stats(trim_dir, sample):
    """Parse TrimGalore statistics files for both reads."""
    stats = {}
    for read in [1, 2]:
        sf = os.path.join(trim_dir, sample, f"trimming_statistics_{read}.txt")
        if not os.path.isfile(sf):
            continue
        with open(sf) as f:
            for ln in f:
                ln = ln.strip()
                if "Total reads processed:" in ln:
                    stats[f"total_reads_R{read}"] = ln.split(":")[-1].strip().replace(",", "")
                elif "Reads with adapters:" in ln:
                    stats[f"adapters_R{read}"] = ln.split(":")[-1].strip().split("(")[0].strip().replace(",", "")
                elif "Reads written (passing filters):" in ln:
                    stats[f"passed_R{read}"] = ln.split(":")[-1].strip().split("(")[0].strip().replace(",", "")
                elif "Quality Phred score cutoff:" in ln:
                    stats["quality_cutoff"] = ln.split(":")[-1].strip()
    return stats if stats else None


def load_bowtie2(log_dir, sample):
    lf = os.path.join(log_dir, sample, "bowtie2_align.log")
    if not os.path.isfile(lf):
        return None
    info = {"overall_rate": None, "total_reads": 0, "paired_pct": 0,
            "concordant_0": 0, "concordant_1": 0, "concordant_multi": 0}
    with open(lf) as f:
        for ln in f:
            ln = ln.strip()
            if "reads; of these:" in ln:
                try: info["total_reads"] = int(ln.split()[0].replace(",", ""))
                except ValueError: pass
            elif "were paired" in ln:
                try: info["paired_pct"] = float(ln.split("(")[1].split("%")[0])
                except (IndexError, ValueError): pass
            elif "aligned concordantly 0 times" in ln and "%" in ln:
                try: info["concordant_0"] = float(ln.split("(")[1].split("%")[0])
                except (IndexError, ValueError): pass
            elif "aligned concordantly exactly 1 time" in ln:
                try: info["concordant_1"] = float(ln.split("(")[1].split("%")[0])
                except (IndexError, ValueError): pass
            elif "aligned concordantly >1 times" in ln:
                try: info["concordant_multi"] = float(ln.split("(")[1].split("%")[0])
                except (IndexError, ValueError): pass
            elif "overall alignment rate" in ln:
                try: info["overall_rate"] = float(ln.split("%")[0].split()[-1])
                except (IndexError, ValueError): pass
    return info


def load_bowtie2_metrics(metrics_dir, sample):
    mf = os.path.join(metrics_dir, sample, f"{sample}_bowtie2_metrics.txt")
    if not os.path.isfile(mf):
        return None
    with open(mf) as f:
        header = f.readline().strip().split("\t")
        vals = f.readline().strip().split("\t")
        if header and vals and len(header) == len(vals):
            return dict(zip(header, vals))
    return None


def load_markdup(metrics_dir, sample):
    mf = os.path.join(metrics_dir, sample, f"{sample}.Markdup-metrics.txt")
    if not os.path.isfile(mf):
        return None
    with open(mf) as f:
        for ln in f:
            ln = ln.strip()
            if ln.startswith("## METRICS CLASS"):
                hdr = next(f, "").strip().split("\t")
                vals = next(f, "").strip().split("\t")
                if hdr and vals and len(hdr) == len(vals):
                    d = dict(zip(hdr, vals))
                    return {
                        "unpaired": int(d.get("UNPAIRED_READS_EXAMINED", 0)),
                        "read_pairs": int(d.get("READ_PAIRS_EXAMINED", 0)),
                        "secondary_supp": int(d.get("SECONDARY_OR_SUPPLEMENTARY_RDS", 0)),
                        "unmapped": int(d.get("UNMAPPED_READS", 0)),
                        "unpaired_dups": int(d.get("UNPAIRED_READ_DUPLICATES", 0)),
                        "dup_pairs": int(d.get("READ_PAIR_DUPLICATES", 0)),
                        "optical_dups": int(d.get("READ_PAIR_OPTICAL_DUPLICATES", 0)),
                        "dup_rate": float(d.get("PERCENT_DUPLICATION", 0)) * 100,
                        "est_lib_size": int(d.get("ESTIMATED_LIBRARY_SIZE", 0)),
                    }
    return None


def load_macs3_info(log_dir, sample):
    lf = os.path.join(log_dir, sample, "macs3.log")
    if not os.path.isfile(lf):
        return None
    info = {"tags_treatment": 0, "tags_control": 0, "fragment_length": 0}
    with open(lf) as f:
        for ln in f:
            if "tags after filtering in treatment" in ln:
                try: info["tags_treatment"] = int(ln.split(":")[-1].strip().split()[0])
                except (ValueError, IndexError): pass
            elif "tags after filtering in control" in ln:
                try: info["tags_control"] = int(ln.split(":")[-1].strip().split()[0])
                except (ValueError, IndexError): pass
            elif "fragment length" in ln and "is" in ln:
                try: info["fragment_length"] = int(ln.split("is")[-1].strip().split()[0])
                except (ValueError, IndexError): pass
    return info


def load_narrow_peaks(peaks_dir, sample, max_rows=5000):
    pf = os.path.join(peaks_dir, sample, f"{sample}_peaks.narrowPeak")
    if not os.path.isfile(pf):
        return []
    peaks = []
    with open(pf) as f:
        for ln in f:
            if ln.startswith("#"): continue
            parts = ln.strip().split("\t")
            if len(parts) >= 10:
                try:
                    peaks.append({
                        "chr": parts[0], "start": int(parts[1]), "end": int(parts[2]),
                        "name": parts[3], "score": int(parts[4]),
                        "signalValue": float(parts[6]), "pvalue": float(parts[7]),
                        "qvalue": float(parts[8]), "peak_pos": int(parts[9]),
                    })
                except (ValueError, IndexError): pass
            if len(peaks) >= max_rows: break
    return peaks


def load_broad_peaks(peaks_dir, sample, max_rows=5000):
    pf = os.path.join(peaks_dir, sample, f"{sample}_broad_peaks.broadPeak")
    if not os.path.isfile(pf):
        return []
    peaks = []
    with open(pf) as f:
        for ln in f:
            if ln.startswith("#"): continue
            parts = ln.strip().split("\t")
            if len(parts) >= 9:
                try:
                    peaks.append({
                        "chr": parts[0], "start": int(parts[1]), "end": int(parts[2]),
                        "name": parts[3], "score": int(parts[4]),
                        "signalValue": float(parts[6]), "pvalue": float(parts[7]),
                        "qvalue": float(parts[8]),
                    })
                except (ValueError, IndexError): pass
            if len(peaks) >= max_rows: break
    return peaks


def load_macs3_xls(peaks_dir, sample, max_rows=5000):
    xf = os.path.join(peaks_dir, sample, f"{sample}_peaks.xls")
    if not os.path.isfile(xf):
        return []
    rows = []
    with open(xf) as f:
        for ln in f:
            if ln.startswith("#") or not ln.strip(): continue
            parts = ln.strip().split("\t")
            if len(parts) >= 10 and parts[0] != "chr":
                try:
                    rows.append({
                        "chr": parts[0], "start": int(parts[1]), "end": int(parts[2]),
                        "length": int(parts[3]), "abs_summit": int(parts[4]),
                        "pileup": int(parts[5]), "neg_log10_pvalue": float(parts[6]),
                        "fold_enrichment": float(parts[7]), "neg_log10_qvalue": float(parts[8]),
                        "name": parts[9],
                    })
                except (ValueError, IndexError): pass
            if len(rows) >= max_rows: break
    return rows


def load_broad_xls(peaks_dir, sample, max_rows=5000):
    xf = os.path.join(peaks_dir, sample, f"{sample}_broad_peaks.xls")
    if not os.path.isfile(xf):
        return []
    rows = []
    with open(xf) as f:
        for ln in f:
            if ln.startswith("#") or not ln.strip(): continue
            parts = ln.strip().split("\t")
            if len(parts) >= 9 and parts[0] != "chr":
                try:
                    rows.append({
                        "chr": parts[0], "start": int(parts[1]), "end": int(parts[2]),
                        "length": int(parts[3]),
                        "pileup": float(parts[4]), "neg_log10_pvalue": float(parts[5]),
                        "fold_enrichment": float(parts[6]), "neg_log10_qvalue": float(parts[7]),
                        "name": parts[8],
                    })
                except (ValueError, IndexError): pass
            if len(rows) >= max_rows: break
    return rows


def load_summits(peaks_dir, sample):
    sf = os.path.join(peaks_dir, sample, f"{sample}_summits.bed")
    if not os.path.isfile(sf):
        return []
    rows = []
    with open(sf) as f:
        for ln in f:
            parts = ln.strip().split("\t")
            if len(parts) >= 5:
                try:
                    rows.append({"chr": parts[0], "start": int(parts[1]), "end": int(parts[2]),
                                 "name": parts[3], "score": float(parts[4])})
                except (ValueError, IndexError): pass
    return rows


def load_cutoff(peaks_dir, sample):
    xf = os.path.join(peaks_dir, sample, f"{sample}_cutoff_analysis.txt")
    if not os.path.isfile(xf):
        return []
    data = []
    with open(xf) as f:
        f.readline()  # skip header
        for ln in f:
            parts = ln.strip().split("\t")
            if len(parts) >= 5:
                try:
                    data.append({"pscore": float(parts[0]), "qscore": float(parts[1]),
                                 "npeaks": int(parts[2]), "lpeaks": int(parts[3]),
                                 "avelpeak": float(parts[4])})
                except ValueError: pass
    return data


def load_broad_cutoff(peaks_dir, sample):
    xf = os.path.join(peaks_dir, sample, f"{sample}_broad_cutoff_analysis.txt")
    if not os.path.isfile(xf):
        return []
    data = []
    with open(xf) as f:
        f.readline()  # skip header
        for ln in f:
            parts = ln.strip().split("\t")
            if len(parts) >= 5:
                try:
                    data.append({"pscore": float(parts[0]), "qscore": float(parts[1]),
                                 "npeaks": int(parts[2]), "lpeaks": int(parts[3]),
                                 "avelpeak": float(parts[4])})
                except ValueError: pass
    return data


def load_frip(qc_dir, sample):
    ff = os.path.join(qc_dir, sample, f"{sample}.FRiP.txt")
    if not os.path.isfile(ff):
        return None
    with open(ff) as f:
        for ln in f:
            parts = ln.strip().split()
            if len(parts) >= 2:
                return float(parts[1])
    return None


def load_peak_count(peaks_dir, sample):
    counts = {"narrow": 0, "broad": 0}
    for suffix, key in [("_peaks.narrowPeak", "narrow"), ("_broad_peaks.broadPeak", "broad")]:
        pf = os.path.join(peaks_dir, sample, f"{sample}{suffix}")
        if os.path.isfile(pf):
            with open(pf) as f:
                counts[key] = sum(1 for ln in f if ln.strip() and not ln.startswith("#"))
    return counts


def load_annotation(annotation_dir, sample):
    af = os.path.join(annotation_dir, sample, f"{sample}_peaks.annotatePeaks.txt")
    if not os.path.isfile(af):
        return None, []
    regions = Counter()
    rows = []
    with open(af) as f:
        reader = csv.reader(f, delimiter="\t")
        header = next(reader, None)
        for row in reader:
            if len(row) >= 10:
                ann = row[7].strip() if len(row) > 7 else ""
                if "Promoter" in ann: regions["promoter"] += 1
                elif "Exon" in ann: regions["exon"] += 1
                elif "Intron" in ann: regions["intron"] += 1
                elif "Intergenic" in ann: regions["intergenic"] += 1
                elif "TTS" in ann: regions["tts"] += 1
                else: regions["other"] += 1
                rows.append({
                    "peak_id": row[0], "chr": row[1], "start": row[2], "end": row[3],
                    "strand": row[4] if len(row) > 4 else "", "score": row[5] if len(row) > 5 else "",
                    "annotation": row[7] if len(row) > 7 else "",
                    "detailed_annotation": row[8] if len(row) > 8 else "",
                    "distance_to_tss": row[9] if len(row) > 9 else "",
                    "nearest_promoter": row[10] if len(row) > 10 else "",
                    "entrez_id": row[11] if len(row) > 11 else "",
                    "gene_name": row[15] if len(row) > 15 else "",
                    "gene_alias": row[16] if len(row) > 16 else "",
                    "gene_description": row[17] if len(row) > 17 else "",
                    "gene_type": row[18] if len(row) > 18 else "",
                })
    return (dict(regions) if regions else None), rows


def load_top_genes(annotation_dir, sample, top_n=5):
    _, rows = load_annotation(annotation_dir, sample)
    scored = [{"gene": r["gene_name"] or r["annotation"],
               "position": f"{r['chr']}:{r['start']}-{r['end']}",
               "score": float(r["score"]) if r["score"] else 0,
               "distance_to_tss": r["distance_to_tss"]}
              for r in rows if r["score"]]
    scored.sort(key=lambda x: -x["score"])
    return scored[:top_n]


def load_tss_distances(annotation_dir, sample):
    _, rows = load_annotation(annotation_dir, sample)
    distances = []
    for r in rows:
        try: distances.append(abs(int(r["distance_to_tss"])))
        except (ValueError, TypeError): pass
    if not distances: return None
    total = len(distances)
    w1 = sum(1 for d in distances if d <= 1000)
    w5 = sum(1 for d in distances if d <= 5000)
    w10 = sum(1 for d in distances if d <= 10000)
    return {"total": total, "distances": distances,
            "within_1kb": w1, "within_1kb_pct": w1/total*100,
            "within_5kb": w5, "within_5kb_pct": w5/total*100,
            "within_10kb": w10, "within_10kb_pct": w10/total*100}


def load_te_overlap_counts(te_dir, sample):
    tf = os.path.join(te_dir, sample, f"{sample}_te_overlap_counts.tsv")
    if not os.path.isfile(tf):
        return []
    rows = []
    with open(tf) as f:
        reader = csv.DictReader(f, delimiter="\t")
        for row in reader:
            rows.append(row)
    return rows


def load_te_subfamily(te_dir, sample):
    tf = os.path.join(te_dir, sample, f"{sample}_te_subfamily_overlap.tsv")
    if not os.path.isfile(tf):
        return []
    rows = []
    with open(tf) as f:
        reader = csv.DictReader(f, delimiter="\t")
        for row in reader:
            rows.append(row)
    return rows


def load_peak_centric_te(te_dir, sample):
    tf = os.path.join(te_dir, sample, f"{sample}_peak_centric_te.tsv")
    if not os.path.isfile(tf):
        return []
    rows = []
    with open(tf) as f:
        reader = csv.DictReader(f, delimiter="\t")
        for row in reader:
            rows.append(row)
    return rows


def get_te_enrichment_png(te_dir, sample, input_sample):
    """Return path to TE enrichment figure if it exists."""
    path = os.path.join(te_dir, f"{sample}_vs_{input_sample}_enrichment.png")
    return path if os.path.isfile(path) else None


def load_fastqc_summary(qc_dir, sample, stage="raw"):
    """Parse FastQC summary txt (raw or trimmed)."""
    suffix = "raw" if stage == "raw" else "trimmed"
    sf = os.path.join(qc_dir, f"{stage}_fastqc", sample, f"{sample}.fastqc.{suffix}.txt")
    if not os.path.isfile(sf):
        return None
    results = {}
    with open(sf) as f:
        for ln in f:
            parts = ln.strip().split("\t")
            if len(parts) >= 2:
                results[parts[0]] = parts[1]  # module_name -> PASS/WARN/FAIL
    return results if results else None


# ============================================================
# Plotting
# ============================================================

def plot_alignment(aligns):
    fig, ax = plt.subplots(figsize=(8, 4))
    samples = list(aligns.keys())
    rates = [aligns[s]["overall_rate"] or 0 for s in samples]
    colors = ["#0094D8" if r >= 80 else "#D64545" for r in rates]
    bars = ax.barh(samples, rates, color=colors, height=0.6)
    ax.set_xlabel("Alignment Rate (%)"); ax.set_title("Bowtie2 Alignment Rate"); ax.set_xlim(0, 105)
    ax.axvline(80, color="red", linestyle="--", alpha=0.5, label="80% threshold")
    for bar, rate in zip(bars, rates):
        ax.text(bar.get_width() + 1, bar.get_y() + bar.get_height()/2, f"{rate:.1f}%", va="center", fontsize=9)
    ax.legend(fontsize=8); plt.tight_layout()
    path = tempfile.mktemp(suffix=".png"); fig.savefig(path, dpi=DPI, bbox_inches="tight"); plt.close(fig)
    return path


def plot_peak_and_frip(samples, peaks, frips):
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 4))
    peak_counts = [peaks[s]["narrow"] for s in samples]
    frip_vals = [(frips[s] or 0) * 100 for s in samples]
    ax1.bar(samples, peak_counts, color="#0094D8", width=0.5)
    ax1.set_title("Peak Count (narrow)"); ax1.set_ylabel("Number of Peaks")
    for i, v in enumerate(peak_counts): ax1.text(i, v + max(peak_counts)*0.02, str(v), ha="center", fontsize=9)
    colors = ["#00A878" if f >= 20 else "#D64545" for f in frip_vals]
    ax2.bar(samples, frip_vals, color=colors, width=0.5)
    ax2.set_title("FRiP Score"); ax2.set_ylabel("FRiP (%)")
    ax2.axhline(20, color="red", linestyle="--", alpha=0.5, label="20% threshold")
    for i, v in enumerate(frip_vals): ax2.text(i, v + max(frip_vals)*0.02, f"{v:.1f}%", ha="center", fontsize=9)
    ax2.legend(fontsize=8); plt.tight_layout()
    path = tempfile.mktemp(suffix=".png"); fig.savefig(path, dpi=DPI, bbox_inches="tight"); plt.close(fig)
    return path


def plot_annotation(anns, lang="zh"):
    if not anns: return None
    labels = ["promoter", "exon", "intron", "intergenic", "tts", "other"]
    samples = list(anns.keys())
    fig, ax = plt.subplots(figsize=(8, 4))
    x = np.arange(len(samples)); width = 0.12
    for i, label in enumerate(labels):
        vals = [anns[s].get(label, 0) for s in samples]
        ax.bar(x + i * width, vals, width, label=label)
    ax.set_xticks(x + width * 2.5); ax.set_xticklabels(samples)
    ax.set_ylabel("Peak Count"); ax.set_title("Peak Annotation — Genomic Region Distribution")
    ax.legend(fontsize=7, ncol=3); plt.tight_layout()
    path = tempfile.mktemp(suffix=".png"); fig.savefig(path, dpi=DPI, bbox_inches="tight"); plt.close(fig)
    return path


def plot_tss_distance(tss_stats):
    fig, ax = plt.subplots(figsize=(8, 4))
    for sample, ts in tss_stats.items():
        if ts and ts.get("distances"):
            bins = np.arange(0, max(min(ts["distances"]), 10000) + 100, 100)
            ax.hist(ts["distances"], bins=bins, alpha=0.5, label=sample, density=True)
    ax.set_xlabel("Distance to TSS (bp)"); ax.set_ylabel("Density")
    ax.set_title("TSS Distance Distribution"); ax.set_xlim(0, 10000); ax.legend(fontsize=8)
    plt.tight_layout()
    path = tempfile.mktemp(suffix=".png"); fig.savefig(path, dpi=DPI, bbox_inches="tight"); plt.close(fig)
    return path


def plot_dup_rate(all_samples, markdups):
    fig, ax = plt.subplots(figsize=(8, 4))
    samples = [s for s in all_samples if s in markdups]
    rates = [markdups[s]["dup_rate"] for s in samples]
    colors = ["#D64545" if r > 50 else "#0094D8" for r in rates]
    bars = ax.barh(samples, rates, color=colors, height=0.6)
    ax.set_xlabel("Duplication Rate (%)"); ax.set_title("MarkDuplicates — Duplication Rate")
    ax.axvline(50, color="red", linestyle="--", alpha=0.5, label="50% threshold")
    for bar, rate in zip(bars, rates):
        ax.text(bar.get_width() + 0.5, bar.get_y() + bar.get_height()/2, f"{rate:.1f}%", va="center", fontsize=9)
    ax.legend(fontsize=8); plt.tight_layout()
    path = tempfile.mktemp(suffix=".png"); fig.savefig(path, dpi=DPI, bbox_inches="tight"); plt.close(fig)
    return path


def plot_trim_stats(all_samples, trims):
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 4))
    samples = [s for s in all_samples if s in trims]
    # Adapter rate (R1)
    adapter_rates = []
    for s in samples:
        t = trims[s]
        total = int(t.get("total_reads_R1", 0))
        adapters = int(t.get("adapters_R1", 0))
        adapter_rates.append(adapters / total * 100 if total > 0 else 0)
    ax1.bar(samples, adapter_rates, color="#F09A36", width=0.5)
    ax1.set_title("Adapter Content (R1)"); ax1.set_ylabel("Reads with Adapters (%)")
    ax1.tick_params(axis='x', rotation=30)
    for i, v in enumerate(adapter_rates): ax1.text(i, v + max(adapter_rates)*0.02, f"{v:.1f}%", ha="center", fontsize=8)
    # Passed reads (R1)
    passed = [int(trims[s].get("passed_R1", 0)) / 1e6 for s in samples]
    ax2.bar(samples, passed, color="#0094D8", width=0.5)
    ax2.set_title("Passed Reads (R1)"); ax2.set_ylabel("Reads (M)")
    ax2.tick_params(axis='x', rotation=30)
    for i, v in enumerate(passed): ax2.text(i, v + max(passed)*0.02, f"{v:.1f}M", ha="center", fontsize=8)
    plt.tight_layout()
    path = tempfile.mktemp(suffix=".png"); fig.savefig(path, dpi=DPI, bbox_inches="tight"); plt.close(fig)
    return path


def plot_peak_centric_te(all_peak_centric_data):
    """Plot TE coverage fraction distribution + TE count per peak."""
    if not all_peak_centric_data:
        return None
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 4))
    # Left: TE coverage fraction histogram
    fracs = [float(r["te_coverage_frac"]) for r in all_peak_centric_data if r.get("te_coverage_frac")]
    if fracs:
        ax1.hist(fracs, bins=30, color="#0094D8", edgecolor="white", alpha=0.8)
        ax1.set_xlabel("TE Coverage Fraction"); ax1.set_ylabel("Peak Count")
        ax1.set_title("Peak TE Coverage Distribution")
        mean_frac = np.mean(fracs)
        ax1.axvline(mean_frac, color="red", linestyle="--", alpha=0.7, label=f"mean={mean_frac:.3f}")
        ax1.legend(fontsize=8)
    # Right: TE count distribution (capped at 20 for display)
    counts = [int(r["te_count"]) for r in all_peak_centric_data if r.get("te_count")]
    if counts:
        capped = [min(c, 20) for c in counts]
        bins = np.arange(0, 22) - 0.5
        ax2.hist(capped, bins=bins, color="#F09A36", edgecolor="white", alpha=0.8)
        ax2.set_xlabel("TE Count per Peak (capped at 20)"); ax2.set_ylabel("Peak Count")
        ax2.set_title("TE Count Distribution")
        ax2.set_xticks(range(0, 21, 2))
    plt.tight_layout()
    path = tempfile.mktemp(suffix=".png"); fig.savefig(path, dpi=DPI, bbox_inches="tight"); plt.close(fig)
    return path


# ============================================================
# PPT Slides
# ============================================================

def build_title_slide(prs, title, subtitle, date, pipeline, lang):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C_NAVY
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(1.8))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_NAVY; bar.line.fill.background()
    _textbox(slide, MARGIN_L, 0.35, CONTENT_W, 0.6, title, font_size=28, bold=True, color=C_WHITE)
    _textbox(slide, MARGIN_L, 1.0, CONTENT_W, 0.4, subtitle, font_size=14, color=RGBColor(0xB0, 0xC4, 0xDE))
    meta = []
    if pipeline: meta.append(f"{t('pipeline_label', lang)}: {pipeline}")
    if date: meta.append(f"{t('date_label', lang)}: {date}")
    if meta: _textbox(slide, MARGIN_L, 2.2, CONTENT_W, 0.3, "  |  ".join(meta), font_size=10, color=C_MUTED)
    # Accent band
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.8), Inches(8.4), Inches(0.15))
    band.fill.solid(); band.fill.fore_color.rgb = C_ACCENT; band.line.fill.background()


def build_workflow_slide(prs, ip_samples, input_samples, lang):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _bg(slide); _header(slide, t("workflow_title", lang))
    steps = [("FastQC", "QC"), ("TrimGalore", "Trim"), ("Bowtie2", "Align"), ("GATK", "MarkDup"),
             ("MACS3", "Peaks"), ("HOMER", "Annot"), ("bamCov", "Tracks"), ("FRiP", "Enrich")]
    n = len(steps); box_w = 0.95; box_h = 0.7; gap = 0.15
    total_w = n * box_w + (n - 1) * gap; x0 = MARGIN_L + (CONTENT_W - total_w) / 2
    for i, (name, desc) in enumerate(steps):
        x = x0 + i * (box_w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.1), Inches(box_w), Inches(box_h))
        shape.fill.solid(); shape.fill.fore_color.rgb = C_NAVY; shape.line.color.rgb = C_ACCENT
        tf = shape.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = f"{i+1}. {name}"; p.font.size = Pt(9); p.font.bold = True
        p.font.color.rgb = C_WHITE; p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(7)
        p2.font.color.rgb = RGBColor(0xA0, 0xB8, 0xD0); p2.alignment = PP_ALIGN.CENTER
        if i < n - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + box_w), Inches(1.35), Inches(gap), Inches(0.2))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = C_ACCENT; arrow.line.fill.background()
    all_s = [(s, t("ip_type", lang)) for s in ip_samples] + [(s, t("input_type", lang)) for s in input_samples]
    data = [[t("sample", lang), t("type", lang)]] + [[s, st] for s, st in all_s]
    _table(slide, MARGIN_L, 2.2, CONTENT_W, 0.28 * len(data), data)
    bullets = [
        "ChIP-seq 标准流程：质控 → 去接头 → 比对 → 去重 → Peak Calling → 注释 → 富集分析" if lang == "zh" else "Standard ChIP-seq pipeline: QC → Trim → Align → Dedup → Peak Calling → Annotation → Enrichment",
        f"IP 样本: {len(ip_samples)}, Input 样本: {len(input_samples)}" if lang == "zh" else f"IP samples: {len(ip_samples)}, Input samples: {len(input_samples)}",
    ]
    _bullets(slide, MARGIN_L, 4.3, CONTENT_W, 0.8, bullets, font_size=11, color=C_MUTED)


def build_trim_slide(prs, all_samples, trims, img_path, lang):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _bg(slide); _header(slide, t("trim_title", lang))
    _add_img(slide, img_path, MARGIN_L, CONTENT_TOP, 6.0, 3.5)
    data = [["Sample", "Total R1", "Adapters R1", "Passed R1", "Cutoff"]]
    for s in all_samples:
        st = trims.get(s, {})
        data.append([s, st.get("total_reads_R1", "-"), st.get("adapters_R1", "-"),
                      st.get("passed_R1", "-"), st.get("quality_cutoff", "-")])
    _table(slide, 6.5, CONTENT_TOP, 3.2, 0.22 * len(data), data, font_size=10)
    bullets = [
        "TrimGalore 去除 adapter 序列和低质量碱基 (Q < 25)" if lang == "zh" else "TrimGalore removes adapter sequences and low-quality bases (Q < 25)",
        "接头含量过高提示文库构建可能存在问题" if lang == "zh" else "High adapter content may indicate library preparation issues",
    ]
    _bullets(slide, MARGIN_L, 4.3, CONTENT_W, 0.8, bullets, font_size=11, color=C_MUTED)


def build_alignment_slide(prs, aligns, img_path, lang):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _bg(slide); _header(slide, t("alignment_title", lang))
    _add_img(slide, img_path, MARGIN_L, CONTENT_TOP, 6.0, 3.5)
    data = [[t("sample", lang), "Reads", "Align%", "Conc 1x", "Conc >1x"]]
    for s, info in aligns.items():
        data.append([s, f"{info['total_reads']:,}",
                      f"{info['overall_rate']:.1f}%" if info['overall_rate'] else "-",
                      f"{info['concordant_1']:.1f}%", f"{info['concordant_multi']:.1f}%"])
    _table(slide, 6.5, CONTENT_TOP, 3.2, 0.22 * len(data), data, font_size=10)
    bullets = [
        "Bowtie2 比对率 ≥ 80% 为合格，< 80% 需检查数据质量和参考基因组" if lang == "zh" else "Alignment rate ≥ 80% is acceptable; < 80% requires data/genome QC",
        "Concordant 比对率反映 paired-end reads 的一致性" if lang == "zh" else "Concordant alignment rate reflects paired-end read consistency",
    ]
    _bullets(slide, MARGIN_L, 4.3, CONTENT_W, 0.8, bullets, font_size=11, color=C_MUTED)


def build_markdup_slide(prs, markdups, all_samples, img_path, lang):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _bg(slide); _header(slide, t("markdup_title", lang))
    _add_img(slide, img_path, MARGIN_L, CONTENT_TOP, 5.5, 3.5)
    data = [[t("sample", lang), t("read_pairs", lang), t("dup_pairs", lang), t("dup_rate", lang), "Lib Size"]]
    for s in all_samples:
        d = markdups.get(s, {})
        data.append([s, f"{d.get('read_pairs', 0):,}", f"{d.get('dup_pairs', 0):,}",
                      f"{d.get('dup_rate', 0):.1f}%", f"{d.get('est_lib_size', 0):,}"])
    _table(slide, 6.2, CONTENT_TOP, 3.5, 0.22 * len(data), data, font_size=10)
    bullets = [
        "GATK MarkDuplicates 标记 PCR 重复，不移除（MACS3 内部处理）" if lang == "zh" else "GATK MarkDuplicates flags PCR duplicates (kept for MACS3 internal handling)",
        "Dup rate > 50% 提示文库复杂度低，可能影响 Peak Calling 质量" if lang == "zh" else "Dup rate > 50% suggests low library complexity, may affect peak calling",
    ]
    _bullets(slide, MARGIN_L, 4.3, CONTENT_W, 0.8, bullets, font_size=11, color=C_MUTED)


def build_peak_slide(prs, img_path, ip_samples, peaks, frips, macs3_info, lang):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _bg(slide); _header(slide, t("peak_calling_title", lang))
    _add_img(slide, img_path, MARGIN_L, CONTENT_TOP, 6.0, 3.0)
    data = [[t("sample", lang), "Narrow", "Broad", t("frip_score", lang), "Tags"]]
    for s in ip_samples:
        data.append([s, str(peaks[s]["narrow"]), str(peaks[s]["broad"]),
                      f"{(frips[s] or 0)*100:.1f}%", f"{macs3_info.get(s, {}).get('tags_treatment', 0):,}"])
    _table(slide, 6.8, CONTENT_TOP, 2.8, 0.25 * len(data), data, font_size=10)
    bullets = [
        "MACS3 以 Input 为对照进行 Peak Calling，识别 IP 样本的富集区域" if lang == "zh" else "MACS3 calls peaks by comparing IP against Input control",
        "FRiP ≥ 20% 为合格，反映 reads 在 peak 区域的富集程度" if lang == "zh" else "FRiP ≥ 20% is acceptable, measures read enrichment in peaks",
    ]
    _bullets(slide, MARGIN_L, 4.0, CONTENT_W, 0.8, bullets, font_size=11, color=C_MUTED)


def build_annotation_slide(prs, img_path, lang):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _bg(slide); _header(slide, t("annotation_title", lang))
    _add_img(slide, img_path, MARGIN_L, CONTENT_TOP, CONTENT_W, 3.5)
    bullets = [
        "HOMER annotatePeaks.pl 将 peaks 注释到最近的基因和基因组区域" if lang == "zh" else "HOMER annotatePeaks.pl annotates peaks to nearest genes and genomic regions",
        "Promoter 区域的 peak 优先关注，可能直接影响基因表达调控" if lang == "zh" else "Peaks in promoter regions are of primary interest for gene regulation",
    ]
    _bullets(slide, MARGIN_L, 4.3, CONTENT_W, 0.8, bullets, font_size=11, color=C_MUTED)


def build_tss_slide(prs, img_path, tss_stats, lang):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _bg(slide); _header(slide, t("tss_title", lang))
    _add_img(slide, img_path, MARGIN_L, CONTENT_TOP, 6.0, 3.5)
    data = [["Sample", "Total", "≤1kb", "≤5kb", "≤10kb"]]
    for s, ts in tss_stats.items():
        if ts:
            data.append([s, str(ts["total"]), f"{ts['within_1kb']} ({ts['within_1kb_pct']:.0f}%)",
                          f"{ts['within_5kb']} ({ts['within_5kb_pct']:.0f}%)",
                          f"{ts['within_10kb']} ({ts['within_10kb_pct']:.0f}%)"])
    _table(slide, 6.8, CONTENT_TOP, 2.8, 0.25 * len(data), data, font_size=10)
    bullets = [
        "TSS 附近富集表明 ChIP 信号集中在转录起始位点，是典型转录因子结合模式" if lang == "zh" else "Enrichment near TSS indicates typical transcription factor binding pattern",
    ]
    _bullets(slide, MARGIN_L, 4.3, CONTENT_W, 0.6, bullets, font_size=11, color=C_MUTED)


def build_top_genes_slide(prs, genes, lang):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _bg(slide); _header(slide, t("top_genes_title", lang))
    all_genes = []
    for sample, genelist in genes.items():
        for g in genelist:
            all_genes.append([sample, g["gene"], g["position"], f"{g['score']:.0f}", g.get("distance_to_tss", "")])
    if all_genes:
        header = ["Sample", "Gene", "Position", "Score", "Dist to TSS"]
        _table(slide, MARGIN_L, CONTENT_TOP, CONTENT_W, 0.25 * (len(all_genes) + 1), [header] + all_genes, font_size=10)
    bullets = [
        "Top 靶基因按 peak score 排序，score 越高代表富集越显著" if lang == "zh" else "Top target genes ranked by peak score — higher score means stronger enrichment",
    ]
    _bullets(slide, MARGIN_L, 4.5, CONTENT_W, 0.6, bullets, font_size=11, color=C_MUTED)


def build_te_enrichment_slide(prs, te_enrichment_imgs, lang):
    """Show TE enrichment figures (IP vs Input)."""
    if not te_enrichment_imgs: return
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _bg(slide); _header(slide, t("te_enrichment_title", lang))
    n = len(te_enrichment_imgs)
    if n == 1:
        _add_img(slide, te_enrichment_imgs[0], MARGIN_L, CONTENT_TOP, CONTENT_W, 3.5)
    elif n == 2:
        _add_img(slide, te_enrichment_imgs[0], MARGIN_L, CONTENT_TOP, 4.5, 3.5)
        _add_img(slide, te_enrichment_imgs[1], 5.2, CONTENT_TOP, 4.5, 3.5)
    else:
        cols = min(n, 3)
        img_w = (CONTENT_W - 0.2 * (cols - 1)) / cols
        for i, img in enumerate(te_enrichment_imgs[:3]):
            _add_img(slide, img, MARGIN_L + i * (img_w + 0.2), CONTENT_TOP, img_w, 3.5)
    bullets = [
        "TE 富集分析比较 IP 与 Input 中 TE subfamily 的 reads 分布" if lang == "zh" else "TE enrichment compares IP vs Input read distribution across TE subfamilies",
        "log2(IP/Input) > 0 表示该 TE 在 IP 中富集，可能被目标蛋白结合" if lang == "zh" else "log2(IP/Input) > 0 indicates TE enrichment in IP, suggesting protein binding",
    ]
    _bullets(slide, MARGIN_L, 4.5, CONTENT_W, 0.8, bullets, font_size=11, color=C_MUTED)


def build_peak_centric_te_slide(prs, img_path, all_peak_centric_data, ip_samples, lang):
    """Peak-centric TE: coverage distribution + top TE-rich peaks + summary stats."""
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _bg(slide)
    _header(slide, t("peak_centric_te_title", lang))
    _add_img(slide, img_path, MARGIN_L, CONTENT_TOP, 6.0, 3.5)
    # Summary stats
    if all_peak_centric_data:
        total = len(all_peak_centric_data)
        with_te = sum(1 for r in all_peak_centric_data if int(r.get("te_count", 0)) > 0)
        te_counts = [int(r.get("te_count", 0)) for r in all_peak_centric_data]
        avg_count = np.mean(te_counts) if te_counts else 0
        fracs = [float(r.get("te_coverage_frac", 0)) for r in all_peak_centric_data]
        avg_frac = np.mean(fracs) if fracs else 0
        sorted_data = sorted(all_peak_centric_data, key=lambda x: -int(x.get("te_count", 0)))[:5]
        data = [["Peak", "TE#","Coverage", "TE Classes"]]
        for r in sorted_data:
            pid = r.get("peak_id", "")
            if len(pid) > 20: pid = pid[:17] + "..."
            data.append([pid, str(r.get("te_count", 0)),
                          f"{float(r.get('te_coverage_frac', 0)):.1%}",
                          r.get("te_classes", "")[:20]])
        data.append(["", "", "", ""])
        data.append(["Total Peaks", str(total), "", ""])
        data.append(["With TE", f"{with_te} ({with_te/total*100:.0f}%)", "", ""])
        data.append(["Avg TE/Peak", f"{avg_count:.1f}", "", ""])
        data.append(["Avg Coverage", f"{avg_frac:.1%}", "", ""])
        _table(slide, 6.2, CONTENT_TOP, 3.5, 0.2 * len(data), data, font_size=10)
    bullets = [
        "Peak 视角分析：每个 peak 内包含多少 TE、什么类型、覆盖比例" if lang == "zh" else "Peak-centric analysis: TE count, types, and coverage per peak",
        "高 TE 覆盖率的 peak 在实验验证时需考虑 TE 的多拷贝性" if lang == "zh" else "Peaks with high TE coverage need multi-copy consideration for validation",
    ]
    _bullets(slide, MARGIN_L, 4.5, CONTENT_W, 0.8, bullets, font_size=11, color=C_MUTED)


def build_heatmap_slide(prs, sample, heatmap_pngs, lang):
    """One slide per sample showing TE enrichment heatmaps in a grid."""
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _bg(slide)
    _header(slide, f"TE Enrichment Heatmaps — {sample}")
    # Show up to 6 heatmaps in 2x3 grid
    n = min(len(heatmap_pngs), 6)
    cols, rows_grid = 3, 2
    img_w = (CONTENT_W - 0.2 * (cols - 1)) / cols
    img_h = 2.0
    for i in range(n):
        r, c = divmod(i, cols)
        x = MARGIN_L + c * (img_w + 0.2)
        y = CONTENT_TOP + r * (img_h + 0.15)
        _add_img(slide, heatmap_pngs[i], x, y, img_w, img_h)
    remaining = len(heatmap_pngs) - n
    bullets = [f"显示 {n}/{len(heatmap_pngs)} 个 TE 热图"] + ([f"还有 {remaining} 个热图未显示"] if remaining > 0 else [])
    _bullets(slide, MARGIN_L, 4.5, CONTENT_W, 0.8, bullets, font_size=11, color=C_MUTED)


def build_summary_slide(prs, ip_samples, aligns, markdups, peaks, frips, lang):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _bg(slide); _header(slide, t("summary_title", lang))
    data = [[t("sample", lang), "Align%", "Narrow", "Broad", t("frip_score", lang), "Dup%", "Status"]]
    for s in ip_samples:
        rate = aligns.get(s, {}).get("overall_rate") or 0
        frip = frips.get(s, 0) or 0
        n_peaks = peaks[s]["narrow"]
        b_peaks = peaks[s]["broad"]
        dup_rate = markdups.get(s, {}).get("dup_rate", 0)
        ok = rate >= ALIGN_THRESH and frip >= FRIP_THRESH and dup_rate <= 50
        status = "✓ OK" if ok else "⚠"
        data.append([s, f"{rate:.1f}%", str(n_peaks), str(b_peaks), f"{frip*100:.1f}%", f"{dup_rate:.1f}%", status])
    _table(slide, MARGIN_L, CONTENT_TOP, CONTENT_W, 0.28 * len(data), data)
    bullets = [
        f"Align% 阈值 ≥{ALIGN_THRESH}%, FRiP 阈值 ≥{FRIP_THRESH*100}%, Dup% 阈值 ≤50%" if lang == "zh" else f"Thresholds: Align% ≥{ALIGN_THRESH}%, FRiP ≥{FRIP_THRESH*100}%, Dup% ≤50%",
        "✓ = 所有指标合格，⚠ = 存在 QC 异常" if lang == "zh" else "✓ = all pass, ⚠ = QC issues detected",
    ]
    _bullets(slide, MARGIN_L, 4.3, CONTENT_W, 0.8, bullets, font_size=11, color=C_MUTED)


def build_conclusion_slide(prs, ip_samples, aligns, frips, peaks, markdups, macs3_info, lang):
    """Conclusion slide with dark background, key findings, and QC cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = C_NAVY

    _textbox(slide, 0.8, 0.4, 4.0, 0.5, "总结与结论" if lang == "zh" else "Conclusions",
             font_size=28, bold=True, color=C_WHITE)

    # Key findings bullets
    bullets = []
    total_narrow = sum(peaks[s]["narrow"] for s in ip_samples)
    avg_align = np.mean([aligns.get(s, {}).get("overall_rate") or 0 for s in ip_samples])
    avg_frip = np.mean([(frips.get(s, 0) or 0) for s in ip_samples])
    all_ok = all(
        (aligns.get(s, {}).get("overall_rate") or 0) >= ALIGN_THRESH and
        (frips.get(s, 0) or 0) >= FRIP_THRESH and
        (markdups.get(s, {}).get("dup_rate", 0)) <= 50
        for s in ip_samples
    )
    if lang == "zh":
        bullets.append(f"共鉴定 {total_narrow:,} 个 narrow peaks（{len(ip_samples)} 个 IP 样本）")
        bullets.append(f"平均比对率 {avg_align:.1f}%，平均 FRiP {avg_frip*100:.1f}%")
        if all_ok:
            bullets.append("所有样本 QC 指标均合格 ✓")
        else:
            bad = [s for s in ip_samples if (aligns.get(s, {}).get("overall_rate") or 0) < ALIGN_THRESH or (frips.get(s, 0) or 0) < FRIP_THRESH]
            bullets.append(f"⚠ 以下样本 QC 异常: {', '.join(bad)}")
    else:
        bullets.append(f"Identified {total_narrow:,} narrow peaks across {len(ip_samples)} IP samples")
        bullets.append(f"Mean alignment rate {avg_align:.1f}%, mean FRiP {avg_frip*100:.1f}%")
        if all_ok:
            bullets.append("All samples pass QC thresholds ✓")
        else:
            bad = [s for s in ip_samples if (aligns.get(s, {}).get("overall_rate") or 0) < ALIGN_THRESH or (frips.get(s, 0) or 0) < FRIP_THRESH]
            bullets.append(f"⚠ QC issues in: {', '.join(bad)}")
    _bullets(slide, 0.9, 1.1, 8.2, 1.4, bullets, font_size=14, color=C_WHITE)

    # QC cards
    cards = [
        ("样本数" if lang == "zh" else "Samples", str(len(ip_samples))),
        ("总 Peaks" if lang == "zh" else "Total Peaks", f"{total_narrow:,}"),
        ("平均 FRiP" if lang == "zh" else "Mean FRiP", f"{avg_frip*100:.1f}%"),
    ]
    for i, (label, value) in enumerate(cards):
        x = 0.8 + i * 2.75
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.8), Inches(2.45), Inches(0.8))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x25, 0x38, 0x5A)
        card.line.color.rgb = C_ACCENT
        tf = card.text_frame; tf.clear()
        p1 = tf.paragraphs[0]; p1.text = label
        p1.font.size = Pt(12); p1.font.color.rgb = RGBColor(0xC9, 0xD7, 0xE8); p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = value
        p2.font.size = Pt(22); p2.font.bold = True; p2.font.color.rgb = C_WHITE; p2.alignment = PP_ALIGN.CENTER

    # Accent band
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.8), Inches(8.4), Inches(0.12))
    band.fill.solid(); band.fill.fore_color.rgb = C_ACCENT; band.line.fill.background()

    # Per-sample mini table
    mini = [["Sample", "Narrow", "Broad", "FRiP", "Dup%", "Align%", "Status"]]
    for s in ip_samples:
        rate = aligns.get(s, {}).get("overall_rate") or 0
        frip = frips.get(s, 0) or 0
        dup = markdups.get(s, {}).get("dup_rate", 0)
        ok = rate >= ALIGN_THRESH and frip >= FRIP_THRESH and dup <= 50
        mini.append([s, str(peaks[s]["narrow"]), str(peaks[s]["broad"]),
                     f"{frip*100:.1f}%", f"{dup:.1f}%", f"{rate:.1f}%",
                     "✓" if ok else "⚠"])
    _table(slide, 0.8, 4.1, 8.4, 0.22 * len(mini), mini, font_size=10)


# ============================================================
# Excel export — ALL data from ALL modules
# ============================================================

def write_results_excel(output_path, ip_samples, input_samples, peaks_dir, annotation_dir,
                        qc_dir, log_dir, markdup_dir, te_dir, trim_dir, metrics_dir):
    import openpyxl
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()
    _used = set()

    def _sheet(name):
        safe = name[:31]
        for ch in "/\\?*[]:": safe = safe.replace(ch, "_")
        if safe in _used:
            for i in range(1, 100):
                cand = f"{safe[:28]}_{i}"
                if cand not in _used: safe = cand; break
        _used.add(safe); return safe

    def _auto_width(ws, max_w=40):
        for col in ws.columns:
            ml = max((len(str(c.value or "")) for c in col), default=8)
            ws.column_dimensions[get_column_letter(col[0].column)].width = min(ml + 2, max_w)

    def _write_dicts(ws, rows, keys=None):
        if not rows: return
        if keys is None: keys = list(rows[0].keys())
        ws.append(keys)
        for r in rows:
            ws.append([r.get(k, "") for k in keys])
        _auto_width(ws)

    all_samples = ip_samples + input_samples

    # 1. Overview
    ws = wb.active; ws.title = _sheet("Overview")
    for row in [["Category", "Item", "Value"],
                ["Analysis", "n_ip", str(len(ip_samples))],
                ["Analysis", "n_input", str(len(input_samples))],
                ["Analysis", "ip_samples", ", ".join(ip_samples)],
                ["Analysis", "input_samples", ", ".join(input_samples)]]:
        ws.append(row)
    _auto_width(ws)

    # 2. MACS3 Narrow Peaks
    ws2 = wb.create_sheet(_sheet("MACS3 Narrow Peaks"))
    _write_dicts(ws2, [{**r, "sample_id": s} for s in ip_samples for r in load_macs3_xls(peaks_dir, s)],
                 ["sample_id", "chr", "start", "end", "length", "abs_summit", "pileup", "neg_log10_pvalue", "fold_enrichment", "neg_log10_qvalue", "name"])

    # 3. MACS3 Broad Peaks
    ws3 = wb.create_sheet(_sheet("MACS3 Broad Peaks"))
    _write_dicts(ws3, [{**r, "sample_id": s} for s in ip_samples for r in load_broad_xls(peaks_dir, s)],
                 ["sample_id", "chr", "start", "end", "length", "pileup", "neg_log10_pvalue", "fold_enrichment", "neg_log10_qvalue", "name"])

    # 4. MACS3 Peak Summits
    ws4 = wb.create_sheet(_sheet("MACS3 Peak Summits"))
    _write_dicts(ws4, [{**r, "sample_id": s} for s in ip_samples for r in load_summits(peaks_dir, s)],
                 ["sample_id", "chr", "start", "end", "name", "score"])

    # 5. Cutoff Narrow
    ws5 = wb.create_sheet(_sheet("Cutoff Narrow"))
    ws5.append(["Sample", "-log10(p)", "-log10(q)", "Num Peaks", "Total Length", "Avg Peak Length"])
    for s in ip_samples:
        for r in load_cutoff(peaks_dir, s):
            ws5.append([s, r["pscore"], r["qscore"], r["npeaks"], r["lpeaks"], r["avelpeak"]])
    _auto_width(ws5)

    # 6. Cutoff Broad
    ws6 = wb.create_sheet(_sheet("Cutoff Broad"))
    ws6.append(["Sample", "-log10(p)", "-log10(q)", "Num Peaks", "Total Length", "Avg Peak Length"])
    for s in ip_samples:
        for r in load_broad_cutoff(peaks_dir, s):
            ws6.append([s, r["pscore"], r["qscore"], r["npeaks"], r["lpeaks"], r["avelpeak"]])
    _auto_width(ws6)

    # 7. HOMER Annotation
    ws7 = wb.create_sheet(_sheet("HOMER Annotation"))
    ann_keys = ["peak_id", "chr", "start", "end", "strand", "score", "annotation",
                "detailed_annotation", "distance_to_tss", "nearest_promoter", "entrez_id",
                "gene_name", "gene_alias", "gene_description", "gene_type"]
    ws7.append(["sample"] + ann_keys)
    for s in ip_samples:
        _, rows = load_annotation(annotation_dir, s)
        for r in rows:
            ws7.append([s] + [r.get(k, "") for k in ann_keys])
    _auto_width(ws7)

    # 8. Region Distribution
    ws8 = wb.create_sheet(_sheet("Region Distribution"))
    ann_labels = ["promoter", "exon", "intron", "intergenic", "tts", "other"]
    ws8.append(["Sample"] + ann_labels + ["Total"])
    for s in ip_samples:
        ann, _ = load_annotation(annotation_dir, s)
        a = ann or {}
        total = sum(a.get(l, 0) for l in ann_labels)
        ws8.append([s] + [a.get(l, 0) for l in ann_labels] + [total])
    _auto_width(ws8)

    # 9. Top Genes
    ws9 = wb.create_sheet(_sheet("Top Genes"))
    ws9.append(["Sample", "Gene", "Position", "Score", "Distance to TSS"])
    for s in ip_samples:
        for g in load_top_genes(annotation_dir, s):
            ws9.append([s, g["gene"], g["position"], g["score"], g["distance_to_tss"]])
    _auto_width(ws9)

    # 10. TSS Distance
    ws10 = wb.create_sheet(_sheet("TSS Distance"))
    ws10.append(["Sample", "Total", "≤1kb", "≤1kb (%)", "≤5kb", "≤5kb (%)", "≤10kb", "≤10kb (%)"])
    for s in ip_samples:
        ts = load_tss_distances(annotation_dir, s)
        if ts:
            ws10.append([s, ts["total"], ts["within_1kb"], round(ts["within_1kb_pct"], 2),
                          ts["within_5kb"], round(ts["within_5kb_pct"], 2),
                          ts["within_10kb"], round(ts["within_10kb_pct"], 2)])
    _auto_width(ws10)

    # 11. TE Subfamily Overlap
    ws11 = wb.create_sheet(_sheet("TE Subfamily Overlap"))
    ws11.append(["Sample", "TE Subfamily", "TE Class", "TE Length", "Interval Overlap Frac", "Overlap Peak Count", "IP Reads", "Input Reads"])
    for s in ip_samples:
        for r in load_te_subfamily(te_dir, s):
            ws11.append([s, r.get("te_subfamily", ""), r.get("te_class", ""), r.get("te_length", ""), r.get("interval_overlap_frac", ""),
                          r.get("overlap_peak_count", ""), r.get("ip_reads", ""), r.get("input_reads", "")])
    _auto_width(ws11)

    # 12. Peak-Centric TE
    ws12 = wb.create_sheet(_sheet("Peak-Centric TE"))
    pc_keys = ["sample_id", "peak_id", "chrom", "start", "end", "peak_length",
               "te_count", "te_class_count", "te_classes",
               "te_subfamily_count", "te_subfamilies",
               "te_covered_bases", "te_coverage_frac"]
    ws12.append(pc_keys)
    for s in ip_samples:
        for r in load_peak_centric_te(te_dir, s):
            ws12.append([r.get(k, "") for k in pc_keys])
    _auto_width(ws12)

    # 13. QC Summary (consolidated, last sheet)
    ws_qc = wb.create_sheet(_sheet("QC Summary"))
    # Build header: core QC columns + Bowtie2 metrics columns
    qc_header = [
        "Sample", "Type",
        "Trim_Total_R1", "Trim_Adapters_R1", "Trim_Passed_R1",
        "Trim_Total_R2", "Trim_Adapters_R2", "Trim_Passed_R2", "Trim_Quality_Cutoff",
        "Align_Total_Reads", "Align_Paired_Pct", "Align_Concordant_0", "Align_Concordant_1",
        "Align_Concordant_GT1", "Align_Overall_Pct",
        "MarkDup_Read_Pairs", "MarkDup_Dup_Pairs", "MarkDup_Dup_Rate", "MarkDup_Optical_Dups",
        "MarkDup_Unmapped", "MarkDup_Est_Lib_Size",
        "MACS3_Tags_Treatment", "MACS3_Tags_Control", "MACS3_Fragment_Length",
        "Peak_Narrow", "Peak_Broad", "FRiP_Score", "FRiP_Pct",
        "Status",
    ]
    # Append Bowtie2 metrics column names
    bt2_keys = []
    for s in all_samples:
        m = load_bowtie2_metrics(metrics_dir, s)
        if m:
            bt2_keys = list(m.keys())
            break
    qc_header += bt2_keys
    ws_qc.append(qc_header)
    for s in all_samples:
        is_ip = s in ip_samples
        st = load_trim_stats(trim_dir, s) or {}
        al = load_bowtie2(log_dir, s) or {}
        md = load_markdup(markdup_dir, s) or {}
        mi = load_macs3_info(log_dir, s) or {} if is_ip else {}
        pc = load_peak_count(peaks_dir, s) if is_ip else {"narrow": "", "broad": ""}
        frip = load_frip(qc_dir, s) if is_ip else None
        # Status (IP only)
        if is_ip:
            rate = al.get("overall_rate") or 0
            dup = md.get("dup_rate", 0)
            issues = []
            if rate < ALIGN_THRESH: issues.append("low_align")
            if (frip or 0) < FRIP_THRESH: issues.append("low_frip")
            if dup > 50: issues.append("high_dup")
            status = "OK" if not issues else "; ".join(issues)
        else:
            status = ""
        row = [
            s, "IP" if is_ip else "Input",
            st.get("total_reads_R1", ""), st.get("adapters_R1", ""), st.get("passed_R1", ""),
            st.get("total_reads_R2", ""), st.get("adapters_R2", ""), st.get("passed_R2", ""),
            st.get("quality_cutoff", ""),
            al.get("total_reads", ""), al.get("paired_pct", ""),
            al.get("concordant_0", ""), al.get("concordant_1", ""),
            al.get("concordant_multi", ""), al.get("overall_rate", ""),
            md.get("read_pairs", ""), md.get("dup_pairs", ""), md.get("dup_rate", ""),
            md.get("optical_dups", ""), md.get("unmapped", ""), md.get("est_lib_size", ""),
            mi.get("tags_treatment", ""), mi.get("tags_control", ""), mi.get("fragment_length", ""),
            pc.get("narrow", ""), pc.get("broad", ""),
            frip, round((frip or 0) * 100, 2) if frip else "",
            status,
        ]
        # Append Bowtie2 metrics values
        bt2 = load_bowtie2_metrics(metrics_dir, s) or {}
        row += [bt2.get(k, "") for k in bt2_keys]
        ws_qc.append(row)
    _auto_width(ws_qc)

    wb.save(output_path)
    print(f"Excel saved: {output_path}")


# ============================================================
# Main
# ============================================================
def main():
    ap = argparse.ArgumentParser(description="Generate ChIP-seq Peak Calling PPT report")
    ap.add_argument("--samples", action="append", default=[], required=True, help="IP sample names (repeatable)")
    ap.add_argument("--input-samples", action="append", default=[], required=True, help="Input sample names (repeatable)")
    ap.add_argument("--ip-input-pair", action="append", default=[], help="IP:Input pair mapping (repeatable)")
    ap.add_argument("--peaks-dir", required=True)
    ap.add_argument("--annotation-dir", required=True)
    ap.add_argument("--qc-dir", required=True)
    ap.add_argument("--log-dir", required=True)
    ap.add_argument("--markdup-dir", required=True)
    ap.add_argument("--trim-dir", default="")
    ap.add_argument("--metrics-dir", default="")
    ap.add_argument("--te-dir", default="")
    ap.add_argument("--heatmap-dir", default="")
    ap.add_argument("--output", required=True)
    ap.add_argument("--excel-output", default="")
    ap.add_argument("--img-dir", default="")
    ap.add_argument("--title", default="")
    ap.add_argument("--subtitle", default="")
    ap.add_argument("--pipeline", default="")
    ap.add_argument("--genome", default="")
    ap.add_argument("--date", default="")
    ap.add_argument("--top-n", type=int, default=5)
    ap.add_argument("--lang", choices=["zh", "en"], default="zh")
    args = ap.parse_args()

    ip_samples = args.samples
    input_samples = args.input_samples
    all_samples = ip_samples + input_samples
    lang = args.lang

    # Parse IP-input pairs
    sample_ip_input_map = {}
    for pair in args.ip_input_pair:
        if ":" in pair:
            ip, inp = pair.split(":", 1)
            sample_ip_input_map[ip] = inp

    if not args.title: args.title = t("report_title", lang)
    if not args.subtitle: args.subtitle = " & ".join(ip_samples) + (f" — {args.genome}" if args.genome else "")

    # Derived dirs
    base_dir = os.path.dirname(args.peaks_dir)  # results/
    te_dir = args.te_dir or os.path.join(base_dir, "te_overlap")
    heatmap_dir = args.heatmap_dir or os.path.join(base_dir, "heatmap")
    trim_dir = args.trim_dir or os.path.join(os.path.dirname(base_dir), "common/2_trimmed_fastq")
    metrics_dir = args.metrics_dir or os.path.join(os.path.dirname(base_dir), "common/3_raw_bam")

    # Load data for PPT
    aligns = {s: load_bowtie2(args.log_dir, s) or {} for s in all_samples}
    markdups = {s: d for s in all_samples if (d := load_markdup(args.markdup_dir, s))}
    macs3_info = {s: d for s in ip_samples if (d := load_macs3_info(args.log_dir, s))}
    peaks = {s: load_peak_count(args.peaks_dir, s) for s in ip_samples}
    frips = {s: load_frip(args.qc_dir, s) for s in ip_samples}
    trims = {s: d for s in all_samples if (d := load_trim_stats(trim_dir, s))}
    anns = {}
    for s in ip_samples:
        reg, _ = load_annotation(args.annotation_dir, s)
        if reg: anns[s] = reg
    genes = {s: load_top_genes(args.annotation_dir, s, args.top_n) for s in ip_samples}
    tss_stats = {s: ts for s in ip_samples if (ts := load_tss_distances(args.annotation_dir, s))}
    te_enrichment_imgs = []
    for ip, inp in sample_ip_input_map.items():
        img = get_te_enrichment_png(te_dir, ip, inp)
        if img:
            te_enrichment_imgs.append(img)
    # Heatmap images per sample
    heatmap_imgs = {}
    for s in ip_samples:
        sample_heatmap_dir = os.path.join(heatmap_dir, s)
        if os.path.isdir(sample_heatmap_dir):
            pngs = sorted([os.path.join(sample_heatmap_dir, f) for f in os.listdir(sample_heatmap_dir) if f.endswith("_heatmap.png")])
            if pngs:
                heatmap_imgs[s] = pngs
    # Peak-centric TE data
    all_peak_centric_data = []
    for s in ip_samples:
        all_peak_centric_data.extend(load_peak_centric_te(te_dir, s))

    # Generate plots
    if args.img_dir: os.makedirs(args.img_dir, exist_ok=True)
    align_img = plot_alignment(aligns)
    pf_img = plot_peak_and_frip(ip_samples, peaks, frips)
    ann_img = plot_annotation(anns, lang) if anns else None
    tss_img = plot_tss_distance(tss_stats) if tss_stats else None
    dup_img = plot_dup_rate(all_samples, markdups)
    trim_img = plot_trim_stats(all_samples, trims) if trims else None
    peak_centric_img = plot_peak_centric_te(all_peak_centric_data) if all_peak_centric_data else None

    all_imgs = {"alignment": align_img, "peak_frip": pf_img, "annotation": ann_img,
                "tss_distance": tss_img, "dup_rate": dup_img, "trim_stats": trim_img,
                "peak_centric_te": peak_centric_img}
    if args.img_dir:
        import shutil
        for name, path in all_imgs.items():
            if path and os.path.exists(path):
                shutil.copy2(path, os.path.join(args.img_dir, f"{name}.png"))

    # Build PPT
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W); prs.slide_height = Inches(SLIDE_H)
    build_title_slide(prs, args.title, args.subtitle, args.date, args.pipeline, lang)
    build_workflow_slide(prs, ip_samples, input_samples, lang)
    if trim_img:
        build_trim_slide(prs, all_samples, trims, trim_img, lang)
    build_alignment_slide(prs, aligns, align_img, lang)
    build_markdup_slide(prs, markdups, all_samples, dup_img, lang)
    build_peak_slide(prs, pf_img, ip_samples, peaks, frips, macs3_info, lang)
    if ann_img: build_annotation_slide(prs, ann_img, lang)
    if tss_img: build_tss_slide(prs, tss_img, tss_stats, lang)
    if any(genes.values()): build_top_genes_slide(prs, genes, lang)
    if te_enrichment_imgs: build_te_enrichment_slide(prs, te_enrichment_imgs, lang)
    if peak_centric_img and all_peak_centric_data: build_peak_centric_te_slide(prs, peak_centric_img, all_peak_centric_data, ip_samples, lang)
    for s, pngs in heatmap_imgs.items():
        build_heatmap_slide(prs, s, pngs, lang)
    build_summary_slide(prs, ip_samples, aligns, markdups, peaks, frips, lang)
    build_conclusion_slide(prs, ip_samples, aligns, frips, peaks, markdups, macs3_info, lang)

    os.makedirs(os.path.dirname(os.path.abspath(args.output)), exist_ok=True)
    prs.save(args.output)
    print(f"Saved: {args.output} ({len(prs.slides)} slides)")

    # Save Excel
    excel_path = args.excel_output or os.path.splitext(args.output)[0] + ".xlsx"
    os.makedirs(os.path.dirname(os.path.abspath(excel_path)), exist_ok=True)
    write_results_excel(excel_path, ip_samples, input_samples,
                        args.peaks_dir, args.annotation_dir, args.qc_dir,
                        args.log_dir, args.markdup_dir, te_dir, trim_dir, metrics_dir)

    if not args.img_dir:
        for img in all_imgs.values():
            if img and os.path.exists(img): os.unlink(img)


if __name__ == "__main__":
    main()
