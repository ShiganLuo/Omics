#!/usr/bin/env python3
import argparse
import math
import os
import shutil
import sys
import tempfile
from collections import OrderedDict
from pathlib import Path
from typing import Dict, List, Optional, Set

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

# ── Chinese font support for matplotlib ──────────────────────────────────
for _font_name in ["WenQuanYi Micro Hei", "WenQuanYi Zen Hei", "Noto Sans CJK SC",
                   "SimHei", "Microsoft YaHei", "Arial Unicode MS"]:
    try:
        matplotlib.font_manager.findfont(_font_name, fallback_to_default=False)
        plt.rcParams["font.sans-serif"] = [_font_name]
        plt.rcParams["axes.unicode_minus"] = False
        break
    except Exception:
        continue

# ── Import venn.py from src/common/plot/Python ───────────────────────────
_VENN_DIR = os.path.join(os.path.dirname(__file__), "..", "..", "..", "src", "common", "plot", "Python")
if os.path.isdir(_VENN_DIR) and _VENN_DIR not in sys.path:
    sys.path.insert(0, os.path.abspath(_VENN_DIR))
try:
    import venn as venn_lib
except ImportError:
    venn_lib = None

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image as PILImage

SLIDE_W = 10.0
SLIDE_H = 5.625
HEADER_H = 0.65
MARGIN_L = 0.45
MARGIN_R = 0.45
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R
CONTENT_TOP = HEADER_H + 0.18
CONTENT_H = SLIDE_H - CONTENT_TOP - 0.22
DPI = 300

C_NAVY = RGBColor(0x18, 0x25, 0x43)
C_ACCENT = RGBColor(0x00, 0x94, 0xD8)
C_GREEN = RGBColor(0x00, 0xA8, 0x78)
C_RED = RGBColor(0xD6, 0x45, 0x45)
C_ORANGE = RGBColor(0xF0, 0x9A, 0x36)
C_TEXT = RGBColor(0x33, 0x33, 0x33)
C_MUTED = RGBColor(0x74, 0x74, 0x74)
C_BG = RGBColor(0xF6, 0xF8, 0xFB)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

I18N = {
    "zh": {
        "title": "RNA-seq 分析报告",
        "date": "报告日期",
        "pipeline": "分析流程",
        "workflow": "流程概览",
        "sample_overview": "样本与转录本概览",
        "te_summary": "TE 嵌合转录本概览",
        "te_detail": "TE 类型分布",
        "de_summary": "差异表达概览",
        "de_detail": "差异表达图谱",
        "fusion_summary": "融合转录本概览",
        "conclusion": "结果总结",
        "sample": "样本",
        "group": "分组",
        "layout": "建库",
        "total_tx": "总转录本",
        "chimeric_tx": "TE嵌合转录本",
        "chim_ratio": "嵌合占比",
        "contrast": "比较",
        "gene": "Gene",
        "te": "TE",
        "gene_te": "Gene+TE",
        "up": "上调",
        "down": "下调",
        "significant": "显著",
        "notes": "要点",
        "fusion_total": "总融合数",
        "fusion_high": "高可信",
        "fusion_support": "最大支持",
        "shared_recurrent": "高频复现融合",
        "support": "支持度",
        "n_samples": "样本数",
        "function_summary": "功能富集概览",
        "function_go_kegg": "GO/KEGG 富集",
        "function_gsea": "GSEA 富集",
        "go": "GO",
        "kegg": "KEGG",
        "up_genes": "上调基因",
        "down_genes": "下调基因",
        "pathway": "通路",
        "nes": "NES",
        "padj": "padj",
        "count": "基因数",
    },
    "en": {
        "title": "RNA-seq Analysis Report",
        "date": "Date",
        "pipeline": "Pipeline",
        "workflow": "Workflow Overview",
        "sample_overview": "Sample and Transcript Overview",
        "te_summary": "TE-Chimeric Transcript Overview",
        "te_detail": "TE Type Distribution",
        "de_summary": "Differential Expression Overview",
        "de_detail": "Differential Expression Plots",
        "fusion_summary": "Fusion Transcript Overview",
        "conclusion": "Conclusions",
        "sample": "Sample",
        "group": "Group",
        "layout": "Layout",
        "total_tx": "Total Transcripts",
        "chimeric_tx": "TE-Chimeric",
        "chim_ratio": "Chimeric Ratio",
        "contrast": "Contrast",
        "gene": "Gene",
        "te": "TE",
        "gene_te": "Gene+TE",
        "up": "Up",
        "down": "Down",
        "significant": "Significant",
        "notes": "Highlights",
        "fusion_total": "Total Fusions",
        "fusion_high": "High Confidence",
        "fusion_support": "Max Support",
        "shared_recurrent": "Recurrent Fusions",
        "support": "Support",
        "n_samples": "Samples",
        "function_summary": "Functional Enrichment Overview",
        "function_go_kegg": "GO/KEGG Enrichment",
        "function_gsea": "GSEA Enrichment",
        "go": "GO",
        "kegg": "KEGG",
        "up_genes": "Up Genes",
        "down_genes": "Down Genes",
        "pathway": "Pathway",
        "nes": "NES",
        "padj": "padj",
        "count": "Count",
    },
}


def t(key: str, lang: str) -> str:
    return I18N.get(lang, I18N["zh"]).get(key, key)


class TempImageStore:
    def __init__(self, img_dir: str = ""):
        self.img_dir = img_dir
        self.paths = []
        if self.img_dir:
            os.makedirs(self.img_dir, exist_ok=True)

    def save_fig(self, fig, stem: str) -> str:
        fd, tmp_path = tempfile.mkstemp(prefix=f"{stem}_", suffix=".png")
        os.close(fd)
        fig.savefig(tmp_path, dpi=DPI, bbox_inches="tight")
        plt.close(fig)
        self.paths.append(tmp_path)
        if self.img_dir:
            dest = os.path.join(self.img_dir, f"{stem}.png")
            shutil.copy2(tmp_path, dest)
            return dest
        return tmp_path

    def cleanup(self):
        for path in self.paths:
            if os.path.exists(path):
                os.unlink(path)


def safe_read_tsv(path: str) -> pd.DataFrame:
    if not path or not os.path.isfile(path):
        return pd.DataFrame()
    return pd.read_csv(path, sep="\t")


def load_group_map(contrast_dirs: list[str]) -> dict[str, str]:
    group_map: dict[str, str] = {}
    for contrast_dir in contrast_dirs:
        group_path = os.path.join(contrast_dir, "group.tsv")
        df = safe_read_tsv(group_path)
        if df.empty or "sample" not in df.columns or "group" not in df.columns:
            continue
        for _, row in df.iterrows():
            group_map[str(row["sample"])] = str(row["group"])
    return group_map


def load_te_sample_summary(analysis_dir: str) -> pd.DataFrame:
    path = os.path.join(analysis_dir, "transcripts", "TE_chimeric", "TE_chimeric_sample_summary.tsv")
    df = safe_read_tsv(path)
    if df.empty:
        return df
    for col in ["total_tx", "chimeric_tx", "five_any", "three_any", "end_any", "internal_any", "none"]:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)
    df["chimeric_ratio"] = df["chimeric_tx"] / df["total_tx"].replace(0, pd.NA)
    return df


def load_te_group_summary(analysis_dir: str) -> pd.DataFrame:
    path = os.path.join(analysis_dir, "transcripts", "TE_chimeric", "TE_chimeric_group_summary.tsv")
    df = safe_read_tsv(path)
    if df.empty:
        return df
    for col in ["five_any", "three_any", "end_any", "internal_any"]:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)
    return df


def load_te_type_counts(analysis_dir: str) -> pd.DataFrame:
    path = os.path.join(analysis_dir, "transcripts", "TE_chimeric", "TE_chimeric_te_type_counts.tsv")
    df = safe_read_tsv(path)
    if df.empty:
        return df
    for col in ["five_count", "three_count", "total_count"]:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)
    return df


def load_fusion_summary(analysis_dir: str) -> pd.DataFrame:
    path = os.path.join(analysis_dir, "fusion", "arriba_report", "per_sample_summary.tsv")
    df = safe_read_tsv(path)
    if df.empty:
        return df
    numeric_cols = [
        "total_fusions", "unique_gene_pairs", "high_confidence", "medium_confidence",
        "low_confidence", "in_frame", "out_of_frame", "avg_support", "max_support"
    ]
    for col in numeric_cols:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)
    return df


def load_recurrent_fusions(analysis_dir: str) -> pd.DataFrame:
    path = os.path.join(analysis_dir, "fusion", "arriba_report", "recurrent_fusions.tsv")
    return safe_read_tsv(path)


def safe_read_csv(path: str) -> pd.DataFrame:
    if not path or not os.path.isfile(path):
        return pd.DataFrame()
    return pd.read_csv(path)


def _count_lines(path: str) -> int:
    if not path or not os.path.isfile(path):
        return 0
    with open(path, "r") as fh:
        return sum(1 for _ in fh)


def load_function_summary(contrast_dir: str) -> dict:
    """Load GO/KEGG/GSEA enrichment results for a single contrast.

    Expects *contrast_dir* to be the function output directory for one contrast,
    e.g. ``{analysis_dir}/function/{contrast_name}``.
    """
    name = os.path.basename(contrast_dir)
    result: dict = {
        "contrast": name,
        "func_dir": contrast_dir,
        "go_plot": os.path.join(contrast_dir, "go_back_to_back.png"),
        "kegg_plot": os.path.join(contrast_dir, "kegg_back_to_back.png"),
        "go_up_path": os.path.join(contrast_dir, "go_up.csv"),
        "go_down_path": os.path.join(contrast_dir, "go_down.csv"),
        "kegg_up_path": os.path.join(contrast_dir, "kegg_up.csv"),
        "kegg_down_path": os.path.join(contrast_dir, "kegg_down.csv"),
        "up_genes_path": os.path.join(contrast_dir, "up_genes.txt"),
        "down_genes_path": os.path.join(contrast_dir, "down_genes.txt"),
        "gsea_plot": os.path.join(contrast_dir, "GSEA", "TEcount_Gene_GSEA.jpeg"),
        "gsea_csv_path": os.path.join(contrast_dir, "GSEA", "TEcount_Gene_GSEA.csv"),
    }
    result["up_count"] = _count_lines(result["up_genes_path"])
    result["down_count"] = _count_lines(result["down_genes_path"])

    for key in ["go_up", "go_down", "kegg_up", "kegg_down"]:
        df = safe_read_csv(result[f"{key}_path"])
        result[f"{key}_df"] = df
        result[f"{key}_n"] = int(len(df.index))

    gsea_df = safe_read_csv(result["gsea_csv_path"])
    if not gsea_df.empty and "NES" in gsea_df.columns:
        gsea_df["NES"] = pd.to_numeric(gsea_df["NES"], errors="coerce")
        gsea_df["padj"] = pd.to_numeric(gsea_df.get("padj"), errors="coerce")
    result["gsea_df"] = gsea_df
    result["gsea_n"] = int(len(gsea_df.index))
    return result


def load_diff_summary(contrast_dir: str) -> dict:
    name = os.path.basename(contrast_dir)
    prefix = "{}.".format(name)
    result = {
        "contrast": name,
        "group_path": os.path.join(contrast_dir, "group.tsv"),
        "pca": os.path.join(contrast_dir, "PCA", "{}cpmPCA.png".format(prefix)),
        "gene_volcano": os.path.join(contrast_dir, "volcano", "{}TEcount_Gene_volcano.png".format(prefix)),
        "gene_heatmap": os.path.join(contrast_dir, "heatmap", "{}TEcount_Gene_updown.png".format(prefix)),
        "gene_updown_path": os.path.join(contrast_dir, "upDown", "{}TEcount_Gene_updown.tsv".format(prefix)),
        "te_updown_path": os.path.join(contrast_dir, "upDown", "{}TEcount_TE_updown.tsv".format(prefix)),
        "gene_te_updown_path": os.path.join(contrast_dir, "upDown", "{}TEcount_Gene_TE_updown.tsv".format(prefix)),
        "gene_name_tsv_path": os.path.join(contrast_dir, "{}TEcount_Gene.name.tsv".format(prefix)),
        "te_name_tsv_path": os.path.join(contrast_dir, "{}TEcount_TE.name.tsv".format(prefix)),
    }
    for key in ["gene", "te", "gene_te"]:
        df = safe_read_tsv(result[f"{key}_updown_path"])
        sig_series = df.get("sig", pd.Series(dtype=str)).astype(str) if not df.empty else pd.Series(dtype=str)
        result[f"{key}_total"] = int(len(df.index)) if not df.empty else 0
        result[f"{key}_up"] = int((sig_series == "up").sum()) if not df.empty else 0
        result[f"{key}_down"] = int((sig_series == "down").sum()) if not df.empty else 0
        if not df.empty and "log2FoldChange" in df.columns:
            top_df = df.copy()
            top_df["log2FoldChange"] = pd.to_numeric(top_df["log2FoldChange"], errors="coerce")
            top_df["padj"] = pd.to_numeric(top_df.get("padj"), errors="coerce")
            top_df = top_df.dropna(subset=["log2FoldChange"]).copy()
            if top_df.shape[0] > 0:
                top_df["abs_lfc"] = top_df["log2FoldChange"].abs()
                top_df = top_df.sort_values(["abs_lfc", "padj"], ascending=[False, True])
                result[f"top_{key}"] = top_df.head(5)
            else:
                result[f"top_{key}"] = pd.DataFrame()
        else:
            result[f"top_{key}"] = pd.DataFrame()
    return result


def _add_picture(slide, path: str, left: float, top: float, max_w: float, max_h: float):
    if not path or not os.path.isfile(path):
        return None
    img = PILImage.open(path)
    aspect = img.size[0] / max(img.size[1], 1)
    width = max_w
    height = width / aspect
    if height > max_h:
        height = max_h
        width = height * aspect
    x = left + (max_w - width) / 2
    slide.shapes.add_picture(path, Inches(x), Inches(top), Inches(width), Inches(height))
    return width, height


def _header(slide, text: str):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(HEADER_H))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_NAVY
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(MARGIN_L), Inches(0.08), Inches(CONTENT_W), Inches(0.45))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C_WHITE


def _textbox(slide, left, top, width, height, text, font_size=11, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def _bullets(slide, left, top, width, height, items, font_size=11, color=C_TEXT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return box


def _table(slide, left, top, width, height, data, font_size=10):
    rows = len(data)
    cols = len(data[0]) if data else 1
    shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    tbl = shape.table
    for r, row in enumerate(data):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(value)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_ACCENT
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xFA)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.alignment = PP_ALIGN.CENTER
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = C_WHITE
                else:
                    p.font.color.rgb = C_TEXT
    return tbl


def plot_chimeric_ratio(df: pd.DataFrame, img_store: TempImageStore) -> str:
    fig, ax1 = plt.subplots(figsize=(7.5, 4.0))
    order = df["sample"].tolist()
    ratios = (df["chimeric_ratio"].fillna(0) * 100).tolist()
    totals = df["chimeric_tx"].tolist()
    bars = ax1.bar(order, ratios, color="#1f77b4", alpha=0.85)
    ax1.set_ylabel("Chimeric Ratio (%)")
    ax1.set_ylim(0, max(100, max(ratios) * 1.15 if ratios else 100))
    ax1.set_title("TE-chimeric Transcript Ratio by Sample", fontsize=13, weight="bold")
    ax1.tick_params(axis="x", rotation=35)
    ax1.grid(axis="y", alpha=0.25)
    ax2 = ax1.twinx()
    ax2.plot(order, totals, color="#d64545", marker="o", linewidth=2)
    ax2.set_ylabel("Chimeric Transcript Count")
    for bar, ratio in zip(bars, ratios):
        ax1.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.5, f"{ratio:.1f}%", ha="center", va="bottom", fontsize=9)
    fig.tight_layout()
    return img_store.save_fig(fig, "chimeric_ratio")


def plot_de_counts(diff_summaries: list[dict], img_store: TempImageStore) -> str:
    labels = [d["contrast"] for d in diff_summaries]
    gene_vals = [d["gene_total"] for d in diff_summaries]
    te_vals = [d["te_total"] for d in diff_summaries]
    gene_te_vals = [d["gene_te_total"] for d in diff_summaries]
    x = list(range(len(labels)))
    width = 0.22
    fig, ax = plt.subplots(figsize=(7.4, 4.0))
    ax.bar([i - width for i in x], gene_vals, width=width, label="Gene", color="#1f77b4")
    ax.bar(x, te_vals, width=width, label="TE", color="#2ca02c")
    ax.bar([i + width for i in x], gene_te_vals, width=width, label="Gene+TE", color="#ff7f0e")
    ax.set_xticks(x)
    ax.set_xticklabels(labels, rotation=20)
    ax.set_ylabel("Significant Features")
    ax.set_title("Differential Expression Summary", fontsize=13, weight="bold")
    ax.grid(axis="y", alpha=0.25)
    ax.legend(frameon=False)
    fig.tight_layout()
    return img_store.save_fig(fig, "de_summary")


# ── Venn / UpSet plot helpers ────────────────────────────────────────────

def load_updown_sets(name_tsv_path: str) -> Dict[str, Set[str]]:
    """Load up/down gene name sets from a TEcount_*.name.tsv file.

    Uses the same thresholds as DESeq2.r ScreenFeature:
    padj < 0.05 and |log2FoldChange| >= 0.58 (lfc_cut default).
    Returns ``{"up": set, "down": set}`` using gene_name (first column).
    """
    result: Dict[str, Set[str]] = {"up": set(), "down": set()}
    df = safe_read_tsv(name_tsv_path)
    if df.empty:
        return result
    name_col = df.columns[0]
    if "log2FoldChange" not in df.columns or "padj" not in df.columns:
        return result
    lfc = pd.to_numeric(df["log2FoldChange"], errors="coerce")
    padj = pd.to_numeric(df["padj"], errors="coerce")
    sig = (padj < 0.05) & (lfc.abs() >= 0.58)
    result["up"] = set(df.loc[sig & (lfc > 0), name_col].astype(str))
    result["down"] = set(df.loc[sig & (lfc < 0), name_col].astype(str))
    return result


def _plot_upset(sets: List[Set[str]], labels: List[str], title: str,
                img_store: "TempImageStore", stem: str) -> str:
    """UpSet plot for 7+ sets (pairwise intersection bar chart)."""
    import itertools
    n = len(sets)
    pairs = list(itertools.combinations(range(n), 2))
    pair_labels = ["{} & {}".format(labels[i], labels[j]) for i, j in pairs]
    inter_sizes = [len(sets[i] & sets[j]) for i, j in pairs]
    x = list(range(len(pairs)))
    fig, ax = plt.subplots(figsize=(max(8, len(pairs) * 0.6), 5))
    ax.bar(x, inter_sizes, color="#4C72B0", edgecolor="white", linewidth=0.5)
    ax.set_xticks(x)
    ax.set_xticklabels(pair_labels, rotation=35, ha="right", fontsize=9)
    ax.set_ylabel("Intersection Size")
    ax.set_title(title, fontsize=13, weight="bold", pad=12)
    for i, s in enumerate(inter_sizes):
        ax.text(i, s + 0.3, str(s), ha="center", va="bottom", fontsize=9)
    ax.grid(axis="y", alpha=0.2)
    fig.tight_layout()
    return img_store.save_fig(fig, stem)


def plot_venn(sets: List[Set[str]], labels: List[str], title: str,
              img_store: "TempImageStore", stem: str) -> str:
    """Plot a Venn diagram (2-6 sets) or UpSet plot (7+).

    Uses venn.py from src/common/plot/Python for 2-6 sets.
    Falls back to UpSet-style bar chart for 7+ sets.
    """
    n = len(sets)

    if n <= 6 and venn_lib is not None:
        # venn.py expects list[Iterable], returns labels dict + (fig, ax)
        labels_dict = venn_lib.get_labels(sets, fill=["number"])
        venn_func = getattr(venn_lib, "venn{}".format(n), None)
        if venn_func is None:
            return _plot_upset(sets, labels, title, img_store, stem)
        # Close any existing figures to avoid venn.py's plt.figure(0) conflicts
        plt.close("all")
        fig, _ax = venn_func(labels_dict, names=labels, figsize=(8, 7), dpi=150)
        fig.suptitle(title, fontsize=14, weight="bold", y=0.98)
        path = img_store.save_fig(fig, stem)
        plt.close("all")  # Clean up venn.py's global figure state
        return path
    else:
        # 7+ sets or venn.py not available: UpSet-style bar chart
        return _plot_upset(sets, labels, title, img_store, stem)


def build_venn_slides(prs: "Presentation", diff_summaries: List[dict], lang: str) -> Dict[str, List[str]]:
    """Build Venn diagram slides for gene/TE up/down across contrasts.

    Returns a dict mapping sheet_name -> list[str] of feature IDs shared
    across all contrasts (for Excel output).
    """
    contrasts = [d["contrast"] for d in diff_summaries]
    intersection_data: Dict[str, List[str]] = {}

    if len(contrasts) < 2:
        return intersection_data

    feature_types = [
        ("gene", "gene_name_tsv_path", "up", "Gene Up", "venn_gene_up"),
        ("gene", "gene_name_tsv_path", "down", "Gene Down", "venn_gene_down"),
        ("te", "te_name_tsv_path", "up", "TE Up", "venn_te_up"),
        ("te", "te_name_tsv_path", "down", "TE Down", "venn_te_down"),
    ]

    for _ft, path_key, direction, slide_title, stem in feature_types:
        sets: List[Set[str]] = []
        labels: List[str] = []
        for d in diff_summaries:
            updown = load_updown_sets(d[path_key])
            sets.append(updown[direction])
            labels.append(pretty_contrast(d["contrast"]))

        if all(len(s) == 0 for s in sets):
            continue

        img_store = prs._img_store  # type: ignore[attr-defined]
        venn_img = plot_venn(sets, labels, slide_title, img_store, stem)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = C_BG
        _header(slide, "{}: {}".format(t("de_detail", lang), slide_title))
        _add_picture(slide, venn_img, 1.0, 0.95, 8.0, 3.8)

        # Collect intersection (all contrasts) for Excel
        if len(sets) >= 2:
            shared = sets[0].copy()
            for s in sets[1:]:
                shared &= s
            if shared:
                sheet_name = stem.replace("venn_", "Venn_")
                intersection_data[sheet_name] = sorted(shared)

    return intersection_data


def build_workflow_slide(prs: Presentation, pipeline_text: str, lang: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG
    _header(slide, t("workflow", lang))
    steps = [
        "FASTQ / Meta",
        "TrimGalore / Cutadapt",
        "STAR / HISAT2",
        "TEcount / StringTie / Arriba",
        "DESeq2 + GO/KEGG/GSEA",
    ]
    box_w = 1.55
    gap = 0.2
    start_x = 0.5
    y = 1.5
    colors = [RGBColor(0xE8, 0xF1, 0xFB), RGBColor(0xE9, 0xF7, 0xF1), RGBColor(0xF9, 0xEE, 0xD7), RGBColor(0xF8, 0xE7, 0xEA), RGBColor(0xE9, 0xE8, 0xFA)]
    for idx, step in enumerate(steps):
        x = start_x + idx * (box_w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(0.95))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[idx]
        shape.line.color.rgb = C_ACCENT
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_NAVY
        p.alignment = PP_ALIGN.CENTER
        if idx < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + box_w + 0.03), Inches(y + 0.28), Inches(0.14), Inches(0.34))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = C_ACCENT
            arrow.line.fill.background()
    _textbox(slide, 0.6, 3.0, 8.8, 1.3, pipeline_text, font_size=14, color=C_TEXT)
    bullets = [
        "TEcount 负责基因与转座子联合定量，DESeq2 输出差异表达结果。" if lang == "zh" else "TEcount quantifies genes and transposable elements jointly; DESeq2 performs differential expression.",
        "clusterProfiler 进行 GO/KEGG 富集，fgsea 进行 GSEA 通路分析。" if lang == "zh" else "clusterProfiler performs GO/KEGG enrichment, and fgsea runs GSEA pathway analysis.",
        "最终报告复用现有 PNG / TSV / CSV 结果，不重复计算上游分析。" if lang == "zh" else "The final report reuses existing PNG/TSV/CSV outputs without recomputing upstream analyses.",
    ]
    _bullets(slide, 0.7, 3.7, 8.6, 1.35, bullets, font_size=12)


def build_title_slide(prs: Presentation, title: str, subtitle: str, date: str, pipeline_text: str, lang: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_NAVY
    _textbox(slide, 0.8, 0.85, 8.4, 0.9, title, font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        _textbox(slide, 1.1, 1.75, 7.8, 0.7, subtitle, font_size=16, color=RGBColor(0xD9, 0xE7, 0xF5), align=PP_ALIGN.CENTER)
    meta_lines = []
    if date:
        meta_lines.append(f"{t('date', lang)}: {date}")
    if pipeline_text:
        meta_lines.append(f"{t('pipeline', lang)}: {pipeline_text}")
    _textbox(slide, 1.0, 2.55, 8.0, 1.25, "\n".join(meta_lines), font_size=12, color=RGBColor(0xC8, 0xD6, 0xE5), align=PP_ALIGN.CENTER)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(4.35), Inches(8.3), Inches(0.32))
    band.fill.solid()
    band.fill.fore_color.rgb = C_ACCENT
    band.line.fill.background()


def build_sample_slide(prs: Presentation, sample_df: pd.DataFrame, chimeric_plot: str, notes: list[str], lang: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG
    _header(slide, t("sample_overview", lang))
    table_rows = [[t("sample", lang), t("group", lang), t("layout", lang), t("chim_ratio", lang)]]
    for _, row in sample_df.iterrows():
        table_rows.append([
            row["sample"],
            row.get("group", "NA"),
            row.get("layout", "NA"),
            f"{row.get('chimeric_ratio', 0) * 100:.1f}%" if pd.notna(row.get("chimeric_ratio")) else "NA",
        ])
    _table(slide, 0.45, 0.95, 3.25, 2.55, table_rows, font_size=10)
    _add_picture(slide, chimeric_plot, 3.95, 0.95, 5.0, 2.6)
    _textbox(slide, 0.55, 3.7, 1.0, 0.25, t("notes", lang), font_size=13, bold=True)
    _bullets(slide, 0.55, 4.0, 8.7, 1.1, notes, font_size=11)


def build_te_slide(prs: Presentation, analysis_dir: str, lang: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG
    _header(slide, t("te_summary", lang))
    img1 = os.path.join(analysis_dir, "transcripts", "TE_chimeric", "TE_chimeric_group_stacked.png")
    img2 = os.path.join(analysis_dir, "transcripts", "TE_chimeric", "TE_chimeric_te_type_top.png")
    _add_picture(slide, img1, 0.45, 0.95, 4.35, 3.2)
    _add_picture(slide, img2, 5.0, 0.95, 4.55, 3.2)
    bullets = [
        "左图展示各样本/分组的 5'/3'/internal TE 嵌合事件构成。" if lang == "zh" else "Left panel shows the composition of 5'/3'/internal TE-chimeric events across samples/groups.",
        "右图汇总全局最常见 TE 类型，可用于快速定位驱动重复元件。" if lang == "zh" else "Right panel summarizes the most frequent TE types and highlights dominant repetitive elements.",
    ]
    _bullets(slide, 0.55, 4.35, 8.7, 0.8, bullets, font_size=11)


def build_te_detail_slide(prs: Presentation, analysis_dir: str, te_group_df: pd.DataFrame, te_type_df: pd.DataFrame, lang: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG
    _header(slide, t("te_detail", lang))
    img = os.path.join(analysis_dir, "transcripts", "TE_chimeric", "TE_chimeric_te_type_by_group.png")
    _add_picture(slide, img, 0.45, 0.95, 6.25, 3.45)
    group_rows = [[t("group", lang), "5'", "3'", "Internal"]]
    if not te_group_df.empty:
        for _, row in te_group_df.head(6).iterrows():
            group_rows.append([row.get("group", "NA"), int(row.get("five_any", 0)), int(row.get("three_any", 0)), int(row.get("internal_any", 0))])
    _table(slide, 6.9, 1.0, 2.55, 2.2, group_rows, font_size=10)


def pretty_contrast(name: str) -> str:
    return name.replace("_vs_", " vs ")


def build_de_summary_slide(prs: Presentation, de_plot: str, diff_summaries: list[dict], lang: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG
    _header(slide, t("de_summary", lang))
    _add_picture(slide, de_plot, 0.45, 0.95, 9.05, 2.55)
    rows = [[t("contrast", lang), f"{t('gene', lang)} {t('significant', lang)}", f"{t('te', lang)} {t('significant', lang)}", f"{t('gene_te', lang)} {t('significant', lang)}", f"{t('gene', lang)} {t('up', lang)}/{t('down', lang)}"]]
    for item in diff_summaries:
        rows.append([
            pretty_contrast(item["contrast"]),
            item["gene_total"],
            item["te_total"],
            item["gene_te_total"],
            f"{item['gene_up']}/{item['gene_down']}",
        ])
    _table(slide, 0.55, 3.7, 8.9, 1.0 + 0.3 * len(rows), rows, font_size=10)


def build_de_detail_slides(prs: Presentation, item: dict, lang: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG
    _header(slide, f"{t('de_detail', lang)}: {pretty_contrast(item['contrast'])}")
    _add_picture(slide, item["pca"], 0.45, 0.95, 4.35, 2.55)
    _add_picture(slide, item["gene_volcano"], 5.05, 0.95, 4.45, 2.55)
    notes = [
        "PCA 用于观察组间分离，火山图用于定位显著差异基因。" if lang == "zh" else "PCA captures group separation and the volcano plot highlights significant DE genes.",
        f"Gene: {item['gene_up']} up / {item['gene_down']} down; TE: {item['te_up']} up / {item['te_down']} down.",
    ]
    _bullets(slide, 0.7, 3.85, 8.6, 0.9, notes, font_size=12)

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    slide2.background.fill.solid()
    slide2.background.fill.fore_color.rgb = C_BG
    _header(slide2, f"{t('de_detail', lang)}: {pretty_contrast(item['contrast'])} heatmap")
    _add_picture(slide2, item["gene_heatmap"], 0.45, 0.95, 5.75, 3.6)
    top_df = item.get("top_gene", pd.DataFrame())
    rows = [["Feature", "log2FC", "padj", "sig"]]
    if isinstance(top_df, pd.DataFrame) and not top_df.empty:
        id_col = top_df.columns[0]
        for _, row in top_df.head(6).iterrows():
            rows.append([
                str(row[id_col])[:24],
                f"{float(row['log2FoldChange']):.2f}",
                f"{float(row['padj']):.2e}" if pd.notna(row['padj']) else "NA",
                str(row.get('sig', '')),
            ])
    _table(slide2, 6.35, 1.0, 3.1, 2.35, rows, font_size=9)


def build_fusion_slide(prs: Presentation, analysis_dir: str, fusion_df: pd.DataFrame, recurrent_df: pd.DataFrame, lang: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG
    _header(slide, t("fusion_summary", lang))
    fig1 = os.path.join(analysis_dir, "fusion", "arriba_report", "figures", "fig1_per_sample_counts.png")
    fig5 = os.path.join(analysis_dir, "fusion", "arriba_report", "figures", "fig5_recurrent_heatmap.png")
    _add_picture(slide, fig1, 0.45, 0.95, 4.35, 2.2)
    _add_picture(slide, fig5, 5.0, 0.95, 4.55, 2.2)
    rows = [[t("sample", lang), t("fusion_total", lang), t("fusion_high", lang), t("fusion_support", lang)]]
    if not fusion_df.empty:
        for _, row in fusion_df.iterrows():
            rows.append([row["sample"], int(row["total_fusions"]), int(row["high_confidence"]), int(row["max_support"])])
    _table(slide, 0.45, 3.35, 4.35, 1.5, rows, font_size=9)
    recur_rows = [[t("shared_recurrent", lang), t("n_samples", lang), "confidence"]]
    if not recurrent_df.empty:
        for _, row in recurrent_df.head(5).iterrows():
            recur_rows.append([str(row.get("fusion", ""))[:28], int(row.get("n_samples", 0)), row.get("confidence", "")])
    _table(slide, 5.0, 3.35, 4.55, 1.5, recur_rows, font_size=9)


def _enrichment_table_rows(df: pd.DataFrame, p_col: str, top_n: int = 5) -> list[list]:
    """Build table rows from a GO/KEGG enrichment DataFrame."""
    rows: list[list] = []
    if df.empty or "Description" not in df.columns or p_col not in df.columns:
        return rows
    df = df.copy()
    df[p_col] = pd.to_numeric(df[p_col], errors="coerce")
    df = df.dropna(subset=[p_col]).sort_values(p_col)
    for _, row in df.head(top_n).iterrows():
        desc = str(row["Description"])
        if len(desc) > 40:
            desc = desc[:37] + "..."
        count_val = row.get("Count") if "Count" in df.columns else None
        count = int(count_val) if pd.notna(count_val) else ""
        pval = float(row[p_col])
        rows.append([desc, count, f"{pval:.2e}"])
    return rows


def build_function_summary_slide(prs: Presentation, func_summaries: list[dict], lang: str):
    """Summary slide: one row per contrast with GO/KEGG/GSEA counts."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG
    _header(slide, t("function_summary", lang))
    rows = [[t("contrast", lang), t("up_genes", lang), t("down_genes", lang),
             f"{t('go', lang)} {t('up', lang)}", f"{t('go', lang)} {t('down', lang)}",
             f"{t('kegg', lang)} {t('up', lang)}", f"{t('kegg', lang)} {t('down', lang)}",
             "GSEA"]]
    for item in func_summaries:
        rows.append([
            pretty_contrast(item["contrast"]),
            item["up_count"],
            item["down_count"],
            item["go_up_n"],
            item["go_down_n"],
            item["kegg_up_n"],
            item["kegg_down_n"],
            item["gsea_n"],
        ])
    _table(slide, 0.45, 0.95, 9.1, 0.4 + 0.32 * len(rows), rows, font_size=10)
    notes = [
        "GO/KEGG 基于差异基因 (|log2FC|>=1, padj<=0.05) 进行富集；GSEA 使用全部基因排序。" if lang == "zh"
        else "GO/KEGG enrichment uses DE genes (|log2FC|>=1, padj<=0.05); GSEA uses the full ranked gene list.",
        "上下调计数反映显著富集条目数，GSEA 列为总通路数。" if lang == "zh"
        else "Up/down counts are significant enriched terms; GSEA column is total pathways tested.",
    ]
    _bullets(slide, 0.55, 0.95 + 0.4 + 0.32 * len(rows) + 0.15, 8.8, 1.0, notes, font_size=11)


def build_function_go_kegg_slide(prs: Presentation, item: dict, lang: str):
    """Per-contrast GO/KEGG slide: back-to-back bar plots + top-term tables."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG
    _header(slide, f"{t('function_go_kegg', lang)}: {pretty_contrast(item['contrast'])}")
    _add_picture(slide, item["go_plot"], 0.45, 0.95, 4.35, 2.5)
    _add_picture(slide, item["kegg_plot"], 5.0, 0.95, 4.55, 2.5)

    # GO top terms table
    go_rows = [[f"{t('go', lang)} {t('up', lang)}", t("count", lang), "p.adj"]]
    go_rows.extend(_enrichment_table_rows(item["go_up_df"], "p.adjust", top_n=2))
    go_rows.append([f"{t('go', lang)} {t('down', lang)}", "", ""])
    go_rows.extend(_enrichment_table_rows(item["go_down_df"], "p.adjust", top_n=2))
    _table(slide, 0.45, 3.55, 4.35, 0.22 * len(go_rows), go_rows, font_size=8)

    # KEGG top terms table
    kegg_rows = [[f"{t('kegg', lang)} {t('up', lang)}", t("count", lang), "p.adj"]]
    kegg_rows.extend(_enrichment_table_rows(item["kegg_up_df"], "p.adjust", top_n=2))
    kegg_rows.append([f"{t('kegg', lang)} {t('down', lang)}", "", ""])
    kegg_rows.extend(_enrichment_table_rows(item["kegg_down_df"], "p.adjust", top_n=2))
    _table(slide, 5.0, 3.55, 4.55, 0.22 * len(kegg_rows), kegg_rows, font_size=8)


def build_function_gsea_slide(prs: Presentation, item: dict, lang: str):
    """Per-contrast GSEA slide: waterfall plot + top pathway table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG
    _header(slide, f"{t('function_gsea', lang)}: {pretty_contrast(item['contrast'])}")
    _add_picture(slide, item["gsea_plot"], 0.45, 0.95, 5.75, 3.6)

    gsea_df = item.get("gsea_df", pd.DataFrame())
    rows = [[t("pathway", lang), t("nes", lang), t("padj", lang)]]
    if not gsea_df.empty and "pathway" in gsea_df.columns:
        sig = gsea_df.dropna(subset=["NES"]).copy()
        if not sig.empty:
            sig = sig.sort_values("NES", ascending=False)
            top_positive = sig.head(4)
            top_negative = sig.sort_values("NES", ascending=True).head(4)
            combined = pd.concat([top_positive, top_negative]).drop_duplicates(subset=["pathway"])
            for _, row in combined.head(8).iterrows():
                pw = str(row["pathway"])
                if len(pw) > 30:
                    pw = pw[:27] + "..."
                nes = float(row["NES"]) if pd.notna(row["NES"]) else 0.0
                padj = float(row["padj"]) if pd.notna(row["padj"]) else float("nan")
                rows.append([pw, f"{nes:.2f}", f"{padj:.2e}" if pd.notna(padj) else "NA"])
    _table(slide, 6.35, 1.0, 3.15, 0.3 * len(rows), rows, font_size=9)
    notes = [
        f"{item['up_count']} {t('up_genes', lang)}, {item['down_count']} {t('down_genes', lang)}" if lang == "zh"
        else f"{item['up_count']} {t('up_genes', lang)}, {item['down_count']} {t('down_genes', lang)}",
    ]
    _bullets(slide, 0.55, 4.7, 5.6, 0.5, notes, font_size=10)


def build_conclusion_slide(prs: Presentation, bullets: list[str], stats: dict, lang: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_NAVY
    _textbox(slide, 0.8, 0.75, 4.0, 0.5, t("conclusion", lang), font_size=28, bold=True, color=C_WHITE)
    _bullets(slide, 0.9, 1.45, 8.2, 2.35, bullets, font_size=15, color=C_WHITE)
    cards = [
        ("Samples", str(stats.get("samples", 0))),
        ("Contrasts", str(stats.get("contrasts", 0))),
        ("Max fusions", str(stats.get("max_fusions", 0))),
    ]
    if lang == "zh":
        cards = [
            ("样本数", str(stats.get("samples", 0))),
            ("比较数", str(stats.get("contrasts", 0))),
            ("最大融合数", str(stats.get("max_fusions", 0))),
        ]
    x0 = 0.8
    for i, (label, value) in enumerate(cards):
        x = x0 + i * 2.75
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.0), Inches(2.45), Inches(0.7))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x25, 0x38, 0x5A)
        card.line.color.rgb = C_ACCENT
        tf = card.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = label
        p1.font.size = Pt(11)
        p1.font.color.rgb = RGBColor(0xC9, 0xD7, 0xE8)
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(20)
        p2.font.bold = True
        p2.font.color.rgb = C_WHITE
        p2.alignment = PP_ALIGN.CENTER
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.8), Inches(8.4), Inches(0.15))
    band.fill.solid()
    band.fill.fore_color.rgb = C_ACCENT
    band.line.fill.background()


def collect_conclusion_bullets(sample_df: pd.DataFrame, fusion_df: pd.DataFrame, diff_summaries: list[dict], func_summaries: list[dict], lang: str) -> list[str]:
    bullets = []
    if not sample_df.empty and "group" in sample_df.columns:
        ratio_df = sample_df.dropna(subset=["chimeric_ratio"]).groupby("group", as_index=False)["chimeric_ratio"].mean()
        if ratio_df.shape[0] >= 2:
            ratio_df = ratio_df.sort_values("chimeric_ratio", ascending=False)
            top = ratio_df.iloc[0]
            bottom = ratio_df.iloc[-1]
            bullets.append(
                f"{top['group']} 组平均 TE 嵌合占比最高 ({top['chimeric_ratio'] * 100:.1f}%)，高于 {bottom['group']} ({bottom['chimeric_ratio'] * 100:.1f}%)。" if lang == "zh" else f"{top['group']} shows the highest mean TE-chimeric ratio ({top['chimeric_ratio'] * 100:.1f}%), above {bottom['group']} ({bottom['chimeric_ratio'] * 100:.1f}%)."
            )
    if diff_summaries:
        for item in diff_summaries:
            bullets.append(
                f"{item['contrast']} 检出 {item['gene_total']} 个显著基因、{item['te_total']} 个显著 TE。" if lang == "zh" else f"{item['contrast']} detected {item['gene_total']} significant genes and {item['te_total']} significant TEs."
            )
    if func_summaries:
        for item in func_summaries:
            go_total = item["go_up_n"] + item["go_down_n"]
            kegg_total = item["kegg_up_n"] + item["kegg_down_n"]
            bullets.append(
                f"{item['contrast']} 功能富集：GO {go_total} 条、KEGG {kegg_total} 条、GSEA {item['gsea_n']} 条通路。" if lang == "zh"
                else f"{item['contrast']}: GO {go_total} terms, KEGG {kegg_total} terms, GSEA {item['gsea_n']} pathways."
            )
    if not fusion_df.empty:
        max_row = fusion_df.sort_values("total_fusions", ascending=False).iloc[0]
        bullets.append(
            f"融合负荷最高的样本为 {max_row['sample']}，共 {int(max_row['total_fusions'])} 个候选融合，高可信 {int(max_row['high_confidence'])} 个。" if lang == "zh" else f"{max_row['sample']} has the highest fusion burden with {int(max_row['total_fusions'])} candidate fusions and {int(max_row['high_confidence'])} high-confidence events."
        )
    if not bullets:
        bullets.append("报告已生成，但可用于总结的上游结果不足。" if lang == "zh" else "The report was generated, but upstream result summaries were limited.")
    return bullets


def build_sample_dataframe(samples: list[str], paired_samples: list[str], single_samples: list[str], group_map: dict[str, str], te_sample_df: pd.DataFrame) -> pd.DataFrame:
    records = []
    te_indexed = te_sample_df.set_index("sample") if not te_sample_df.empty and "sample" in te_sample_df.columns else pd.DataFrame()
    for sample in samples:
        layout = "PE" if sample in set(paired_samples) else "SE" if sample in set(single_samples) else "NA"
        row = {"sample": sample, "group": group_map.get(sample, "NA"), "layout": layout, "chimeric_ratio": float("nan")}
        if not te_indexed.empty and sample in te_indexed.index:
            row["chimeric_ratio"] = te_indexed.loc[sample, "chimeric_ratio"]
            row["total_tx"] = te_indexed.loc[sample, "total_tx"]
            row["chimeric_tx"] = te_indexed.loc[sample, "chimeric_tx"]
        records.append(row)
    return pd.DataFrame(records)


def write_file_inventory(output_path: str, analysis_dir: str, contrasts: List[str],
                         samples: List[str], paired_samples: List[str],
                         single_samples: List[str],
                         intersection_data: Optional[Dict[str, List[str]]] = None) -> None:
    """Write all result data to an Excel workbook.

    Each result TSV/CSV file becomes a sheet with its full contents.
    Intersection gene lists from Venn diagrams become additional sheets.
    """
    import openpyxl

    wb = openpyxl.Workbook()

    # ── Sheet 1: Overview ─────────────────────────────────────
    ws = wb.active
    ws.title = "Overview"
    overview_rows = [
        ["Category", "Item", "Value"],
        ["Analysis", "analysis_dir", analysis_dir],
        ["Analysis", "n_samples", str(len(samples))],
        ["Analysis", "n_paired", str(len(paired_samples))],
        ["Analysis", "n_single", str(len(single_samples))],
        ["Analysis", "n_contrasts", str(len(contrasts))],
        ["Analysis", "contrasts", ", ".join(contrasts)],
        ["Analysis", "samples", ", ".join(samples)],
    ]
    for row in overview_rows:
        ws.append(row)
    ws.column_dimensions["A"].width = 12
    ws.column_dimensions["B"].width = 18
    ws.column_dimensions["C"].width = 60

    def _safe_sheet_name(name: str) -> str:
        """Excel sheet name: max 31 chars, no : \\ / ? * [ ]"""
        for ch in ":\\/?*[]":
            name = name.replace(ch, "_")
        return name[:31]

    _used_names: set = set()

    def _unique_sheet_name(name: str) -> str:
        """Ensure sheet name is unique and <= 31 chars, auto-number on conflict."""
        base = _safe_sheet_name(name)
        if base not in _used_names:
            _used_names.add(base)
            return base
        for i in range(1, 100):
            suffix = str(i)
            candidate = base[: 31 - len(suffix)] + suffix
            if candidate not in _used_names:
                _used_names.add(candidate)
                return candidate
        return base

    _MAX_ROWS = 5000

    def _write_tsv_sheet(sheet_name: str, tsv_path: str, contrast: str = ""):
        """Write a TSV file's contents to a new sheet (capped at 5000 rows).

        If *contrast* is given, it is written as the first row before the data.
        """
        df = safe_read_tsv(tsv_path)
        if df.empty:
            return
        if len(df) > _MAX_ROWS:
            df = df.head(_MAX_ROWS)
        ws_data = wb.create_sheet(_unique_sheet_name(sheet_name))
        if contrast:
            ws_data.append([contrast])
        ws_data.append(list(df.columns))
        for _, row in df.iterrows():
            ws_data.append([str(row[col]) if pd.notna(row[col]) else "" for col in df.columns])
        for i, col in enumerate(df.columns, 1):
            max_len = max(len(str(col)), int(df[col].astype(str).str.len().max()) if len(df) > 0 else 0)
            ws_data.column_dimensions[openpyxl.utils.get_column_letter(i)].width = min(max_len + 2, 50)

    def _write_csv_sheet(sheet_name: str, csv_path: str, contrast: str = ""):
        """Write a CSV file's contents to a new sheet (capped at 5000 rows).

        If *contrast* is given, it is written as the first row before the data.
        """
        df = safe_read_csv(csv_path)
        if df.empty:
            return
        if len(df) > _MAX_ROWS:
            df = df.head(_MAX_ROWS)
        ws_data = wb.create_sheet(_unique_sheet_name(sheet_name))
        if contrast:
            ws_data.append([contrast])
        ws_data.append(list(df.columns))
        for _, row in df.iterrows():
            ws_data.append([str(row[col]) if pd.notna(row[col]) else "" for col in df.columns])
        for i, col in enumerate(df.columns, 1):
            max_len = max(len(str(col)), int(df[col].astype(str).str.len().max()) if len(df) > 0 else 0)
            ws_data.column_dimensions[openpyxl.utils.get_column_letter(i)].width = min(max_len + 2, 50)

    def _write_text_sheet(sheet_name: str, txt_path: str, contrast: str = ""):
        """Write a text file (one item per line) to a new sheet.

        If *contrast* is given, it is written as the first row before the data.
        """
        if not os.path.isfile(txt_path):
            return
        ws_data = wb.create_sheet(_unique_sheet_name(sheet_name))
        if contrast:
            ws_data.append([contrast])
        ws_data.append(["Feature"])
        with open(txt_path, encoding="utf-8") as f:
            for line in f:
                ws_data.append([line.strip()])
        ws_data.column_dimensions["A"].width = 30

    # ── TE chimeric results (no contrast) ─────────────────────
    te_base = os.path.join(analysis_dir, "transcripts", "TE_chimeric")
    _write_tsv_sheet("te_sample_summary", os.path.join(te_base, "TE_chimeric_sample_summary.tsv"))
    _write_tsv_sheet("te_group_summary", os.path.join(te_base, "TE_chimeric_group_summary.tsv"))
    _write_tsv_sheet("te_type_counts", os.path.join(te_base, "TE_chimeric_te_type_counts.tsv"))

    # ── TEcount matrix ────────────────────────────────────────
    tecount_path = os.path.join(analysis_dir, "results", "counts", "TEcount", "all_TEcount.tsv")
    if os.path.isfile(tecount_path) and os.path.getsize(tecount_path) < 10 * 1024 * 1024:
        _write_tsv_sheet("tecount_matrix", tecount_path)

    # ── Fusion results (no contrast) ──────────────────────────
    fusion_base = os.path.join(analysis_dir, "fusion", "arriba_report")
    _write_tsv_sheet("fusion_summary", os.path.join(fusion_base, "per_sample_summary.tsv"))
    _write_tsv_sheet("recurrent_fusions", os.path.join(fusion_base, "recurrent_fusions.tsv"))
    _write_tsv_sheet("high_medium_fusions", os.path.join(fusion_base, "high_medium_confidence_fusions.tsv"))
    _write_tsv_sheet("inframe_fusions", os.path.join(fusion_base, "inframe_fusions.tsv"))

    # ── Per-contrast DE results ───────────────────────────────
    for contrast in contrasts:
        de_dir = os.path.join(analysis_dir, "diff_expression", contrast)
        contrast_prefix = "{}.".format(contrast)
        for fname, sheet_name in [
            ("{}TEcount_Gene.name.tsv".format(contrast_prefix), "contrast_gene_name"),
            ("{}TEcount_TE.name.tsv".format(contrast_prefix), "contrast_te_name"),
            ("{}TEcount_Gene_TE.name.tsv".format(contrast_prefix), "contrast_gene_te_name"),
        ]:
            _write_tsv_sheet(sheet_name, os.path.join(de_dir, fname), contrast=contrast)


    # ── Per-contrast function results ─────────────────────────
    for contrast in contrasts:
        func_dir = os.path.join(analysis_dir, "function", contrast)
        _write_csv_sheet("func_go_up", os.path.join(func_dir, "go_up.csv"), contrast=contrast)
        _write_csv_sheet("func_go_down", os.path.join(func_dir, "go_down.csv"), contrast=contrast)
        _write_csv_sheet("func_kegg_up", os.path.join(func_dir, "kegg_up.csv"), contrast=contrast)
        _write_csv_sheet("func_kegg_down", os.path.join(func_dir, "kegg_down.csv"), contrast=contrast)
        _write_text_sheet("func_up_genes", os.path.join(func_dir, "up_genes.txt"), contrast=contrast)
        _write_text_sheet("func_down_genes", os.path.join(func_dir, "down_genes.txt"), contrast=contrast)
        _write_csv_sheet("func_gsea_csv", os.path.join(func_dir, "GSEA", "TEcount_Gene_GSEA.csv"), contrast=contrast)

    # ── Intersection sheets from Venn diagrams ────────────────
    if intersection_data:
        for sheet_name, feature_ids in intersection_data.items():
            ws_inter = wb.create_sheet(_unique_sheet_name(sheet_name))
            ws_inter.append(["Feature_Name"])
            for fid in feature_ids:
                ws_inter.append([fid])
            ws_inter.column_dimensions["A"].width = min(
                max(len(str(f)) for f in feature_ids) + 2, 40)

    wb.save(output_path)


def main():
    ap = argparse.ArgumentParser(description="Generate RNAseq PPT report")
    ap.add_argument("--analysis-dir", required=True)
    ap.add_argument("--output", required=True)
    ap.add_argument("--samples", nargs="*", default=[])
    ap.add_argument("--paired-samples", nargs="*", default=[])
    ap.add_argument("--single-samples", nargs="*", default=[])
    ap.add_argument("--contrasts", nargs="*", default=[])
    ap.add_argument("--title", default="")
    ap.add_argument("--subtitle", default="")
    ap.add_argument("--pipeline", default="")
    ap.add_argument("--genome", default="")
    ap.add_argument("--date", default="")
    ap.add_argument("--lang", default="zh")
    ap.add_argument("--img-dir", default="")
    ap.add_argument("--file-inventory", default="",
                    help="If set, write all input file paths to this Excel file")
    args = ap.parse_args()

    analysis_dir = os.path.abspath(args.analysis_dir)
    output = os.path.abspath(args.output)
    os.makedirs(os.path.dirname(output), exist_ok=True)

    contrast_dirs = [os.path.join(analysis_dir, "diff_expression", c) for c in args.contrasts if os.path.isdir(os.path.join(analysis_dir, "diff_expression", c))]
    if not contrast_dirs and os.path.isdir(os.path.join(analysis_dir, "diff_expression")):
        contrast_dirs = [os.path.join(analysis_dir, "diff_expression", d) for d in sorted(os.listdir(os.path.join(analysis_dir, "diff_expression"))) if os.path.isdir(os.path.join(analysis_dir, "diff_expression", d))]

    group_map = load_group_map(contrast_dirs)
    te_sample_df = load_te_sample_summary(analysis_dir)
    te_group_df = load_te_group_summary(analysis_dir)
    te_type_df = load_te_type_counts(analysis_dir)
    fusion_df = load_fusion_summary(analysis_dir)
    recurrent_df = load_recurrent_fusions(analysis_dir)
    diff_summaries = [load_diff_summary(cdir) for cdir in contrast_dirs]
    sample_df = build_sample_dataframe(args.samples, args.paired_samples, args.single_samples, group_map, te_sample_df)

    func_base_dir = os.path.join(analysis_dir, "function")
    func_contrast_dirs = [os.path.join(func_base_dir, c) for c in args.contrasts if os.path.isdir(os.path.join(func_base_dir, c))]
    if not func_contrast_dirs and os.path.isdir(func_base_dir):
        func_contrast_dirs = [os.path.join(func_base_dir, d) for d in sorted(os.listdir(func_base_dir)) if os.path.isdir(os.path.join(func_base_dir, d))]
    func_summaries = [load_function_summary(cdir) for cdir in func_contrast_dirs]

    pipeline_text = args.pipeline or "FASTQ -> TrimGalore/Cutadapt -> STAR/HISAT2 -> TEcount/StringTie/Arriba -> DESeq2 + GO/KEGG/GSEA -> RNAseq_report"
    title = args.title or t("title", args.lang)
    subtitle_parts = []
    if args.subtitle:
        subtitle_parts.append(args.subtitle)
    if args.genome:
        subtitle_parts.append(args.genome)
    subtitle = " | ".join(subtitle_parts)

    img_store = TempImageStore(args.img_dir)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs._img_store = img_store  # type: ignore[attr-defined]

    venn_intersections: Dict[str, List[str]] = {}

    build_title_slide(prs, title, subtitle, args.date, pipeline_text, args.lang)
    build_workflow_slide(prs, pipeline_text, args.lang)

    if not sample_df.empty and not te_sample_df.empty:
        chimeric_plot = plot_chimeric_ratio(te_sample_df, img_store)
        sample_notes = [
            f"共 {len(sample_df)} 个样本；PE {len(args.paired_samples)} 个，SE {len(args.single_samples)} 个。" if args.lang == "zh" else f"{len(sample_df)} samples total; {len(args.paired_samples)} PE and {len(args.single_samples)} SE.",
            f"TE 嵌合转录本占比范围 {te_sample_df['chimeric_ratio'].min() * 100:.1f}% - {te_sample_df['chimeric_ratio'].max() * 100:.1f}%。" if args.lang == "zh" else f"TE-chimeric ratios range from {te_sample_df['chimeric_ratio'].min() * 100:.1f}% to {te_sample_df['chimeric_ratio'].max() * 100:.1f}%.",
        ]
        build_sample_slide(prs, sample_df, chimeric_plot, sample_notes, args.lang)

    if os.path.isfile(os.path.join(analysis_dir, "transcripts", "TE_chimeric", "TE_chimeric_group_stacked.png")):
        build_te_slide(prs, analysis_dir, args.lang)
    if os.path.isfile(os.path.join(analysis_dir, "transcripts", "TE_chimeric", "TE_chimeric_te_type_by_group.png")):
        build_te_detail_slide(prs, analysis_dir, te_group_df, te_type_df, args.lang)

    if diff_summaries:
        de_plot = plot_de_counts(diff_summaries, img_store)
        build_de_summary_slide(prs, de_plot, diff_summaries, args.lang)
        for item in diff_summaries:
            build_de_detail_slides(prs, item, args.lang)
        venn_intersections = build_venn_slides(prs, diff_summaries, args.lang)

    if not fusion_df.empty:
        build_fusion_slide(prs, analysis_dir, fusion_df, recurrent_df, args.lang)

    if func_summaries:
        build_function_summary_slide(prs, func_summaries, args.lang)
        for item in func_summaries:
            if os.path.isfile(item["go_plot"]) or os.path.isfile(item["kegg_plot"]):
                build_function_go_kegg_slide(prs, item, args.lang)
            if os.path.isfile(item["gsea_plot"]):
                build_function_gsea_slide(prs, item, args.lang)

    conclusion_bullets = collect_conclusion_bullets(sample_df, fusion_df, diff_summaries, func_summaries, args.lang)
    build_conclusion_slide(prs, conclusion_bullets, {
        "samples": len(sample_df),
        "contrasts": len(diff_summaries),
        "max_fusions": int(fusion_df["total_fusions"].max()) if not fusion_df.empty else 0,
    }, args.lang)

    prs.save(output)
    img_store.cleanup()

    if args.file_inventory:
        inventory_path = os.path.abspath(args.file_inventory)
        os.makedirs(os.path.dirname(inventory_path), exist_ok=True)
        write_file_inventory(inventory_path, analysis_dir, args.contrasts,
                             args.samples, args.paired_samples, args.single_samples,
                             venn_intersections)

    print(output)


if __name__ == "__main__":
    main()
