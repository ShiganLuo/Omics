#!/usr/bin/env python3
"""Generate a PPT report for ncRNAseq pipeline results.

Collects QC statistics from:
  1. Raw FASTQ read counts
  2. Demultiplex + trimming statistics
  3. STAR 3-pass alignment logs
  4. Per-gene alignment manifests and Tailer CSV outputs
  5. Small-RNA gene annotation BED

Produces a multi-slide PPTX and an Excel file inventory.
"""

from __future__ import annotations

import argparse
import gzip
import os
import re
import shutil
import sys
import tempfile
from collections import OrderedDict, defaultdict
from pathlib import Path
from typing import Dict, List, Optional, Set, Any, Tuple

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as ticker
import numpy as np
import pandas as pd
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
import openpyxl

# ── Chinese font support for matplotlib ──────────────────────────────────
for _font_name in ["WenQuanYi Micro Hei", "WenQuanYi Zen Hei", "Noto Sans CJK SC",
                   "SimHei", "Microsoft YaHei", "Arial Unicode MS"]:
    try:
        matplotlib.font_manager.findfont(_font_name, fallback_to_default=False)
        plt.rcParams["font.sans-serif"] = [_font_name]
        plt.rcParams["axes.unicode_minus"] = False
        break
    except Exception:
        continue

# ── Constants ────────────────────────────────────────────────────────────
SLIDE_W = 10.0
SLIDE_H = 5.625
HEADER_H = 0.65
MARGIN_L = 0.45
MARGIN_R = 0.45
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R
CONTENT_TOP = HEADER_H + 0.18
CONTENT_H = SLIDE_H - CONTENT_TOP - 0.22
DPI = 300

C_NAVY = RGBColor(0x18, 0x25, 0x43)
C_DEEP_BLUE = RGBColor(0x06, 0x5A, 0x82)
C_TEAL = RGBColor(0x1C, 0x72, 0x93)
C_ACCENT = RGBColor(0x00, 0x94, 0xD8)
C_GREEN = RGBColor(0x00, 0xA8, 0x78)
C_RED = RGBColor(0xD6, 0x45, 0x45)
C_ORANGE = RGBColor(0xF0, 0x9A, 0x36)
C_CORAL = RGBColor(0xE7, 0x6F, 0x51)
C_TEXT = RGBColor(0x33, 0x33, 0x33)
C_MUTED = RGBColor(0x74, 0x74, 0x74)
C_BG = RGBColor(0xF6, 0xF8, 0xFB)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT = RGBColor(0xE2, 0xE8, 0xF0)

# ── Plot colors (hex strings for matplotlib) ─────────────────────────────
P_7SL = "#065A82"
P_U1 = "#1C7293"
P_CORAL = "#E76F51"
P_GREEN = "#2A9D8F"
P_ORANGE = "#F4A261"
P_NAVY = "#264653"

I18N = {
    "zh": {
        "title": "ncRNAseq 分析报告",
        "date": "报告日期",
        "pipeline": "分析流程",
        "workflow": "流程概览",
        "read_flow": "Read 流量统计",
        "trimming": "去重与修剪统计",
        "star_mapping": "STAR 3-pass 比对结果",
        "gene_assignment": "基因特异性比对",
        "tail_analysis": "Tailer 尾部分析",
        "tail_length": "Tail 长度分布",
        "end_position": "3' 端位置分布",
        "conclusion": "结果总结",
        "sample": "样本",
        "group": "分组",
        "layout": "建库",
        "raw_reads": "原始 Reads",
        "dedup_reads": "去重后",
        "trimmed_reads": "修剪后",
        "star_input": "STAR 输入",
        "final_bam": "最终 BAM",
        "gene_bam": "基因 BAM",
        "adapter_r1": "R1 Adapter%",
        "adapter_r2": "R2 Adapter%",
        "pairs_removed": "去除对%",
        "unique": "唯一比对",
        "multi": "多重比对",
        "unmapped": "未比对",
        "pass3a": "pass3a (canonical)",
        "pass3b": "pass3b (chimeric)",
        "genes": "基因数",
        "total_assigned": "总分配",
        "top_gene": "主要基因",
        "top_reads": "主要基因 Reads",
        "total_tail": "Tailer Reads",
        "with_tail": "有 Tail",
        "no_tail": "无 Tail",
        "tail_pct": "Tail 占比",
        "upstream": "上游 (<0)",
        "at_end": "末端 (0)",
        "downstream": "下游 (>0)",
        "notes": "要点",
        "count": "Read 数",
    },
    "en": {
        "title": "ncRNAseq Analysis Report",
        "date": "Date",
        "pipeline": "Pipeline",
        "workflow": "Workflow Overview",
        "read_flow": "Read Flow Statistics",
        "trimming": "Dedup & Trimming Statistics",
        "star_mapping": "STAR 3-pass Alignment Results",
        "gene_assignment": "Gene-specific Alignment",
        "tail_analysis": "Tailer Tail Analysis",
        "tail_length": "Tail Length Distribution",
        "end_position": "3' End Position Distribution",
        "conclusion": "Conclusions",
        "sample": "Sample",
        "group": "Group",
        "layout": "Layout",
        "raw_reads": "Raw Reads",
        "dedup_reads": "Post-dedup",
        "trimmed_reads": "Post-trim",
        "star_input": "STAR Input",
        "final_bam": "Final BAM",
        "gene_bam": "Gene BAM",
        "adapter_r1": "R1 Adapter%",
        "adapter_r2": "R2 Adapter%",
        "pairs_removed": "Pairs Removed%",
        "unique": "Unique",
        "multi": "Multi-mapping",
        "unmapped": "Unmapped",
        "pass3a": "pass3a (canonical)",
        "pass3b": "pass3b (chimeric)",
        "genes": "Genes",
        "total_assigned": "Total Assigned",
        "top_gene": "Top Gene",
        "top_reads": "Top Gene Reads",
        "total_tail": "Tailer Reads",
        "with_tail": "With Tail",
        "no_tail": "No Tail",
        "tail_pct": "Tail %",
        "upstream": "Upstream (<0)",
        "at_end": "At end (0)",
        "downstream": "Downstream (>0)",
        "notes": "Highlights",
        "count": "Read count",
    },
}


def t(key: str, lang: str) -> str:
    """Retrieve internationalized text for a given key and language."""
    return I18N.get(lang, I18N["zh"]).get(key, key)


def _clean_list(values: List[Optional[Any]], default: Any = 0) -> List[Any]:
    """Helper: Replace None in a list with a default value (useful for matplotlib)."""
    return [v if v is not None else default for v in values]


def _safe_add(a: Optional[int], b: Optional[int]) -> Optional[int]:
    """Helper: Safely add two numbers that might be None."""
    if a is None and b is None:
        return None
    return (a or 0) + (b or 0)


def _filter_none(values: List[Optional[Any]]) -> List[Any]:
    """Helper: Filter out None values from a list (useful for min/max/sum)."""
    return [v for v in values if v is not None]


# ══════════════════════════════════════════════════════════════════════════
# Data collection
# ══════════════════════════════════════════════════════════════════════════

def count_fastq_reads(path: str) -> Optional[int]:
    """Count reads in a FASTQ file, supporting symbolic links. Returns None if file is missing."""
    real_path = os.path.realpath(path)

    if not os.path.isfile(real_path):
        return None

    opener = gzip.open if real_path.endswith(".gz") else open
    try:
        with opener(real_path, "rt") as fh:
            return sum(1 for _ in fh) // 4
    except Exception:
        return None


def parse_trimming_stats(path: str) -> dict:
    """Parse trim_galore trimming_statistics_1.txt for adapter% and pairs removed%.
    Returns dictionary with parsed values or None for missing data.
    """
    result = {
        "r1_adapter_counts": None, 
        "r1_adapter_pct": None,
        "r2_adapter_counts": None, 
        "r2_adapter_pct": None,
        "cutoff_threshold": None,
        "pairs_removed_counts": None,
        "pairs_removed_pct": None
    }
    if not os.path.isfile(path):
        return result
    text = Path(path).read_text(encoding="utf-8", errors="replace")
    m = re.search(
            r"Reads with adapters:\s+([\d.]+)\s*\(([\d.]+)%\)",
            text
        )
    if m:
        result["r1_adapter_counts"] = float(m.group(1))
        result["r1_adapter_pct"] = float(m.group(2))
        
    stat2_path = path.replace("_1.txt", "_2.txt")
    if os.path.isfile(stat2_path):
        text2 = Path(stat2_path).read_text(encoding="utf-8", errors="replace")
        m2 = re.search(r"Reads with adapters:\s+([\d.]+)\s*\(([\d.]+)%\)", text2)
        if m2:
            result["r2_adapter_counts"] = float(m2.group(1))
            result["r2_adapter_pct"] = float(m2.group(2))
        m3 = re.search(r"Number of sequence pairs removed.*?\(\s*(\d+).*?\):\s+(\d+)\s*\(([\d.]+)%\)", text2)
        if m3:
            result["cutoff_threshold"] = f"{m3.group(1)} bp"
            result["pairs_removed_counts"] = float(m3.group(2))
            result["pairs_removed_pct"] = float(m3.group(3))
    return result


def parse_star_log(path: str) -> dict:
    """Parse STAR Log.final.out for mapping statistics.
    Returns numeric counts and percentages, or None if missing.
    """
    result = {
        "star_input_reads": None, 
        "unique": None, "unique_pct": None,
        "multi": None, "multi_pct": None, 
        "unmapped": None, "unmapped_pct": None
    }
    if not os.path.isfile(path):
        return result
    
    text = Path(path).read_text(encoding="utf-8", errors="replace")
    
    def _grep(pattern: str) -> Optional[int]:
        m = re.search(pattern, text)
        return int(m.group(1).strip()) if m else None

    input_reads = _grep(r"Number of input reads\s+\|\s+(\d+)")
    unique_reads = _grep(r"Uniquely mapped reads number\s+\|\s+(\d+)")
    multi_reads = _grep(r"Number of reads mapped to multiple loci\s+\|\s+(\d+)")
    
    unmapped_short = _grep(r"Number of reads unmapped: too short\s+\|\s+(\d+)") or 0
    unmapped_mismatch = _grep(r"Number of reads unmapped: too many mismatches\s+\|\s+(\d+)") or 0
    unmapped_other = _grep(r"Number of reads unmapped: other\s+\|\s+(\d+)") or 0
    unmapped_reads = unmapped_short + unmapped_mismatch + unmapped_other if input_reads is not None else None

    result["star_input_reads"] = input_reads
    result["unique"] = unique_reads
    result["multi"] = multi_reads
    result["unmapped"] = unmapped_reads
    
    if input_reads and input_reads > 0:
        if unique_reads is not None:
            result["unique_pct"] = (unique_reads / input_reads) * 100
        if multi_reads is not None:
            result["multi_pct"] = (multi_reads / input_reads) * 100
        if unmapped_reads is not None:
            result["unmapped_pct"] = (unmapped_reads / input_reads) * 100

    return result


def load_gene_manifest(path: str) -> pd.DataFrame:
    """Load per-gene alignment manifest TSV into a DataFrame."""
    if not os.path.isfile(path):
        return pd.DataFrame()
    return pd.read_csv(path, sep="\t")


def load_tail_csv(path: str) -> pd.DataFrame:
    """Load Tailer tail analysis CSV into a DataFrame."""
    if not os.path.isfile(path):
        return pd.DataFrame()
    return pd.read_csv(path, skipinitialspace=True)


def collect_sample_data(analysis_dir: str, samples: List[str]) -> List[dict]:
    """Collect all per-sample statistics across the pipeline steps."""
    results = []
    for s in samples:
        d: dict = {"sample": s}
        # 1. Raw FASTQ
        raw_r1 = os.path.join(analysis_dir, "common", "1_raw_fastq", s, f"{s}_1.fq.gz")
        raw_r2 = os.path.join(analysis_dir, "common", "1_raw_fastq", s, f"{s}_2.fq.gz")
        d["raw_reads"] = _safe_add(count_fastq_reads(raw_r1), count_fastq_reads(raw_r2))
        
        # 2. Dedup reads (jla-demultiplexer output)
        dedup_r1 = os.path.join(analysis_dir, "common", "2_trimmed_dedup_fastq", "final_trimmed_fastq", s, f"{s}_1.fq.gz")
        dedup_r2 = os.path.join(analysis_dir, "common", "2_trimmed_dedup_fastq", "final_trimmed_fastq", s, f"{s}_2.fq.gz")
        d["dedup_reads"] = _safe_add(count_fastq_reads(dedup_r1), count_fastq_reads(dedup_r2))
        
        # 3. Trimming stats & Trimmed reads
        trim_stat_r1 = os.path.join(analysis_dir, "common", "2_trimmed_dedup_fastq", "final_trimmed_fastq", s, "trimming_statistics_1.txt")
        d.update(parse_trimming_stats(trim_stat_r1))
        
        trimmed_r1 = os.path.join(analysis_dir, "common", "2_trimmed_dedup_fastq", "final_trimmed_fastq", s, f"{s}_1.fq.gz")
        trimmed_r2 = os.path.join(analysis_dir, "common", "2_trimmed_dedup_fastq", "final_trimmed_fastq", s, f"{s}_2.fq.gz")
        d["trimmed_reads"] = _safe_add(count_fastq_reads(trimmed_r1), count_fastq_reads(trimmed_r2))
        
        # 4. STAR log
        star_log = os.path.join(analysis_dir, "common", "3_raw_bam", s, "star.Log.final.out")
        d.update(parse_star_log(star_log))

        # 5. Final BAM reads (star_3pass output bam)
        gene_dir = os.path.join(analysis_dir, "common", "4_per_gene_bam", s)
        final_bam_file = os.path.join(gene_dir, f"{s}.bam")
        d["final_bam_reads"] = None
        if os.path.exists(final_bam_file):
            try:
                import subprocess
                r = subprocess.run(["samtools", "view", "-c", final_bam_file],
                                   capture_output=True, text=True, timeout=30)
                if r.returncode == 0:
                    d["final_bam_reads"] = int(r.stdout.strip())
            except Exception:
                pass
                
        # 6. Gene manifest
        manifest_path = os.path.join(gene_dir, "genes.tsv")
        manifest = load_gene_manifest(manifest_path)
        d["n_genes"] = len(manifest) if not manifest.empty else 0
        if not manifest.empty and "assigned_records" in manifest.columns:
            d["total_assigned"] = int(manifest["assigned_records"].sum())
            top_row = manifest.sort_values("assigned_records", ascending=False).iloc[0]
            d["top_gene"] = str(top_row["gene_id"])
            d["top_gene_reads"] = int(top_row["assigned_records"])
        else:
            d["total_assigned"] = None
            d["top_gene"] = None
            d["top_gene_reads"] = None
            
        # 7. Tail CSV
        tail_csv = os.path.join(gene_dir, f"{s}_tail.csv")
        tail_df = load_tail_csv(tail_csv)
        if not tail_df.empty:
            d["total_tail_reads"] = int(tail_df["Count"].sum())
            d["with_tail"] = int(tail_df.loc[tail_df["Tail_Length"] > 0, "Count"].sum())
            d["without_tail"] = d["total_tail_reads"] - d["with_tail"]
            # End position distribution
            d["upstream"] = int(tail_df.loc[tail_df["End_Position"] < 0, "Count"].sum())
            d["at_end"] = int(tail_df.loc[tail_df["End_Position"] == 0, "Count"].sum())
            d["downstream"] = int(tail_df.loc[tail_df["End_Position"] > 0, "Count"].sum())
            # Tail length distribution (for reads with tail)
            tail_with = tail_df[tail_df["Tail_Length"] > 0].copy()
            d["tail_lengths"] = tail_with["Tail_Length"].tolist()
            d["tail_counts"] = tail_with["Count"].tolist()
        else:
            d["total_tail_reads"] = None
            d["with_tail"] = None
            d["without_tail"] = None
            d["upstream"] = None
            d["at_end"] = None
            d["downstream"] = None
            d["tail_lengths"] = []
            d["tail_counts"] = []
            
        results.append(d)
    return results


def infer_group(sample: str) -> str:
    """Infer group label from sample name (e.g. 7sl / U1 / other)."""
    lower = sample.lower()
    if "7sl" in lower:
        return "7SL"
    if "u1" in lower:
        return "U1"
    return "Other"


# ══════════════════════════════════════════════════════════════════════════
# PPT helpers
# ══════════════════════════════════════════════════════════════════════════

class TempImageStore:
    """Stores temporary plot images, copying to optional static dir if specified."""
    def __init__(self, img_dir: str = ""):
        self.img_dir = img_dir
        self.paths: list[str] = []
        if self.img_dir:
            os.makedirs(self.img_dir, exist_ok=True)

    def save_fig(self, fig, stem: str) -> str:
        """Saves a matplotlib figure and returns the path."""
        fd, tmp_path = tempfile.mkstemp(prefix=f"{stem}_", suffix=".png")
        os.close(fd)
        fig.savefig(tmp_path, dpi=DPI, bbox_inches="tight")
        plt.close(fig)
        self.paths.append(tmp_path)
        if self.img_dir:
            dest = os.path.join(self.img_dir, f"{stem}.png")
            shutil.copy2(tmp_path, dest)
            return dest
        return tmp_path

    def cleanup(self):
        """Cleans up temporal files."""
        for path in self.paths:
            if os.path.exists(path):
                os.unlink(path)


def _add_picture(slide, path: str, left: float, top: float, max_w: float, max_h: float):
    if not path or not os.path.isfile(path):
        return None
    img = PILImage.open(path)
    aspect = img.size[0] / max(img.size[1], 1)
    width = max_w
    height = width / aspect
    if height > max_h:
        height = max_h
        width = height * aspect
    x = left + (max_w - width) / 2
    slide.shapes.add_picture(path, Inches(x), Inches(top), Inches(width), Inches(height))
    return width, height


def _header(slide, text: str):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(HEADER_H))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_NAVY
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(MARGIN_L), Inches(0.08), Inches(CONTENT_W), Inches(0.45))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C_WHITE


def _textbox(slide, left, top, width, height, text, font_size=11, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def _bullets(slide, left, top, width, height, items, font_size=11, color=C_TEXT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = f"\u2022 {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return box


def _table(slide, left, top, width, height, data, font_size=10):
    rows = len(data)
    cols = len(data[0]) if data else 1
    shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    tbl = shape.table
    for r, row in enumerate(data):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            # Make sure None translates to an empty string in table
            cell.text = str(value) if value is not None else ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_DEEP_BLUE
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xFA)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.alignment = PP_ALIGN.CENTER
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = C_WHITE
                else:
                    p.font.color.rgb = C_TEXT
    return tbl


def _short_name(sample: str) -> str:
    """Shorten sample name for display (e.g. removing common IP prefixes)."""
    return sample.replace("IP_srp54_control_", "")


# ══════════════════════════════════════════════════════════════════════════
# Chart generators
# ══════════════════════════════════════════════════════════════════════════

def plot_read_flow(samples_data: list[dict], img_store: TempImageStore) -> str:
    """Generate Read Flow bar chart across raw/dedup/trim/STAR levels."""
    names = [_short_name(d["sample"]) for d in samples_data]
    groups = [infer_group(d["sample"]) for d in samples_data]
    fig, ax = plt.subplots(figsize=(8.5, 4.0))
    x = np.arange(len(names))
    w = 0.15
    ax.bar(x - 2*w, _clean_list([d.get("raw_reads") for d in samples_data]), w, label="Raw", color="#A8DADC")
    ax.bar(x - w, _clean_list([d.get("dedup_reads") for d in samples_data]), w, label="Post-dedup", color=P_7SL)
    ax.bar(x, _clean_list([d.get("trimmed_reads") for d in samples_data]), w, label="Post-trim", color=P_U1)
    ax.bar(x + w, _clean_list([d.get("final_bam_reads") for d in samples_data]), w, label="STAR 3-pass BAM", color=P_CORAL)
    ax.set_ylabel("Read count", fontsize=11)
    ax.set_xticks(x)
    ax.set_xticklabels(names, fontsize=8, rotation=30, ha="right")
    ax.yaxis.set_major_formatter(ticker.FuncFormatter(lambda v, _: f"{v/1000:.0f}k" if v >= 1000 else f"{v:.0f}"))
    ax.legend(fontsize=8, loc="upper left", framealpha=0.9)
    ax.grid(axis="y", color="#E0E0E0", linewidth=0.5)
    ax.set_axisbelow(True)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    # Group separator
    for i in range(1, len(groups)):
        if groups[i] != groups[i-1]:
            ax.axvline(x=i - 0.5, color="#999999", linewidth=1, linestyle="--", alpha=0.5)
    fig.tight_layout()
    return img_store.save_fig(fig, "read_flow")


def plot_trimming(samples_data: list[dict], img_store: TempImageStore) -> str:
    """Generate Trimming and adapter statistics bar chart."""
    names = [_short_name(d["sample"]) for d in samples_data]
    groups = [infer_group(d["sample"]) for d in samples_data]
    fig, ax = plt.subplots(figsize=(8.5, 4.0))
    x = np.arange(len(names))
    w = 0.25
    ax.bar(x - w, _clean_list([d.get("r1_adapter_pct") for d in samples_data]), w, label="R1 adapter %", color=P_7SL)
    ax.bar(x, _clean_list([d.get("r2_adapter_pct") for d in samples_data]), w, label="R2 adapter %", color=P_U1)
    ax.bar(x + w, _clean_list([d.get("pairs_removed_pct") for d in samples_data]), w, label="Pairs removed %", color=P_CORAL)
    ax.set_ylabel("Percentage (%)", fontsize=11)
    ax.set_xticks(x)
    ax.set_xticklabels(names, fontsize=8, rotation=30, ha="right")
    ax.legend(fontsize=8, framealpha=0.9)
    ax.grid(axis="y", color="#E0E0E0", linewidth=0.5)
    ax.set_axisbelow(True)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    for i in range(1, len(groups)):
        if groups[i] != groups[i-1]:
            ax.axvline(x=i - 0.5, color="#999999", linewidth=1, linestyle="--", alpha=0.5)
    fig.tight_layout()
    return img_store.save_fig(fig, "trimming")


def plot_star_mapping(samples_data: list[dict], img_store: TempImageStore) -> str:
    """Generate STAR final alignment stats and mapping ratio breakdown chart."""
    names = [_short_name(d["sample"]) for d in samples_data]
    fig, axes = plt.subplots(1, 2, figsize=(9, 4.0))
    # Left: final_bam_reads
    ax = axes[0]
    x = np.arange(len(names))
    ax.bar(x, _clean_list([d.get("final_bam_reads") for d in samples_data]), 0.5, color=P_7SL, label="Total mapped")
    ax.set_title("STAR 3-pass Final BAM", fontsize=12, fontweight="bold")
    ax.set_xticks(x)
    ax.set_xticklabels(names, fontsize=7, rotation=30, ha="right")
    ax.yaxis.set_major_formatter(ticker.FuncFormatter(lambda v, _: f"{v/1000:.0f}k" if v >= 1000 else f"{v:.0f}"))
    ax.grid(axis="y", color="#E0E0E0", linewidth=0.5)
    ax.set_axisbelow(True)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    # Right: mapping breakdown
    ax = axes[1]
    unique = _clean_list([d.get("unique") for d in samples_data])
    multi = _clean_list([d.get("multi") for d in samples_data])
    unmapped = _clean_list([d.get("unmapped") for d in samples_data])
    ax.bar(x, unique, 0.5, label="Unique", color=P_7SL)
    ax.bar(x, multi, 0.5, bottom=unique, label="Multi", color=P_U1)
    ax.bar(x, unmapped, 0.5, bottom=[u+m for u, m in zip(unique, multi)], label="Unmapped", color=P_CORAL)
    ax.set_title("pass1 Mapping Breakdown", fontsize=12, fontweight="bold")
    ax.set_xticks(x)
    ax.set_xticklabels(names, fontsize=7, rotation=30, ha="right")
    ax.legend(fontsize=7, framealpha=0.9)
    ax.grid(axis="y", color="#E0E0E0", linewidth=0.5)
    ax.set_axisbelow(True)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.suptitle("STAR 3-pass Alignment Results", fontsize=13, fontweight="bold", y=1.01)
    fig.tight_layout()
    return img_store.save_fig(fig, "star_mapping")


def plot_gene_assignments(samples_data: list[dict], analysis_dir: str, img_store: TempImageStore) -> str:
    """Generate Stacked bar chart of per-gene read assignments, split by group."""
    groups_data: dict[str, dict[str, list[int]]] = defaultdict(lambda: defaultdict(list))
    sample_labels: dict[str, list[str]] = defaultdict(list)
    for d in samples_data:
        g = infer_group(d["sample"])
        sample_labels[g].append(_short_name(d["sample"]))
        manifest_path = os.path.join(analysis_dir, "common", "4_per_gene_bam", d["sample"], "genes.tsv")
        manifest = load_gene_manifest(manifest_path)
        if not manifest.empty and "assigned_records" in manifest.columns:
            for _, row in manifest.sort_values("assigned_records", ascending=False).iterrows():
                groups_data[g][str(row["gene_id"])].append(int(row["assigned_records"]))
            # Pad with zeros for genes not in this sample
            max_genes = max(len(v) for v in groups_data[g].values()) if groups_data[g] else 0
            for gene in groups_data[g]:
                while len(groups_data[g][gene]) < max_genes:
                    groups_data[g][gene].append(0)
        else:
            groups_data[g]["(none)"].append(0)

    n_groups = len(groups_data)
    fig, axes = plt.subplots(1, max(n_groups, 1), figsize=(5 * max(n_groups, 1), 4.0), squeeze=False)
    axes = axes[0]
    palette = [P_7SL, P_U1, "#3A9BBC", "#5EEAD4", "#7BC4D9", "#B8E0E8", "#D0F0F0", P_CORAL, P_GREEN, P_ORANGE]
    for ax_idx, (group, gene_data) in enumerate(sorted(groups_data.items())):
        ax = axes[ax_idx]
        labels = sample_labels[group]
        n_samples = len(labels)
        bottom = np.zeros(n_samples)
        for i, (gene, counts) in enumerate(sorted(gene_data.items(), key=lambda x: -sum(x[1]))):
            vals = np.array(counts[:n_samples])
            ax.bar(range(n_samples), vals, 0.5, bottom=bottom, label=gene, color=palette[i % len(palette)])
            bottom += vals
        ax.set_title(f"{group} group", fontsize=12, fontweight="bold")
        ax.set_xticks(range(n_samples))
        ax.set_xticklabels(labels, fontsize=9)
        ax.set_ylabel("Assigned reads", fontsize=10)
        ax.legend(fontsize=7, loc="upper right", framealpha=0.9)
        ax.grid(axis="y", color="#E0E0E0", linewidth=0.5)
        ax.set_axisbelow(True)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
    fig.suptitle("Per-Gene Read Assignments", fontsize=13, fontweight="bold", y=1.01)
    fig.tight_layout()
    return img_store.save_fig(fig, "gene_assignments")


def plot_tail_summary(samples_data: list[dict], img_store: TempImageStore) -> str:
    """Generate Stacked bar chart for reads with/without non-template tails."""
    names = [_short_name(d["sample"]) for d in samples_data]
    groups = [infer_group(d["sample"]) for d in samples_data]
    fig, ax = plt.subplots(figsize=(8.5, 4.0))
    x = np.arange(len(names))
    w = 0.5
    without = _clean_list([d.get("without_tail") for d in samples_data])
    with_tail = _clean_list([d.get("with_tail") for d in samples_data])
    ax.bar(x, without, w, label="No tail (tail=0)", color=P_7SL)
    ax.bar(x, with_tail, w, bottom=without, label="With tail (tail>0)", color=P_CORAL)
    ax.set_ylabel("Read count", fontsize=11)
    ax.set_xticks(x)
    ax.set_xticklabels(names, fontsize=8, rotation=30, ha="right")
    ax.legend(fontsize=8, framealpha=0.9)
    ax.grid(axis="y", color="#E0E0E0", linewidth=0.5)
    ax.set_axisbelow(True)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    for i in range(1, len(groups)):
        if groups[i] != groups[i-1]:
            ax.axvline(x=i - 0.5, color="#999999", linewidth=1, linestyle="--", alpha=0.5)
    # Percentage labels
    for i, d in enumerate(samples_data):
        total = d.get("total_tail_reads") or 0
        if total > 0:
            pct = d["with_tail"] / total * 100
            ax.text(i, without[i] + with_tail[i] + max(without) * 0.02,
                    f"{pct:.1f}%", ha="center", fontsize=7, color=P_CORAL, fontweight="bold")
    fig.tight_layout()
    return img_store.save_fig(fig, "tail_summary")


def plot_tail_length(samples_data: list[dict], img_store: TempImageStore) -> str:
    """Generate Tail length distribution chart, split by sample groups."""
    groups_data: dict[str, dict[int, int]] = defaultdict(lambda: defaultdict(int))
    for d in samples_data:
        g = infer_group(d["sample"])
        lengths = d.get("tail_lengths") or []
        counts = d.get("tail_counts") or []
        for tl, cnt in zip(lengths, counts):
            groups_data[g][tl] += cnt

    n_groups = len(groups_data) if groups_data else 1
    fig, axes = plt.subplots(1, max(n_groups, 1), figsize=(5 * max(n_groups, 1), 4.0), squeeze=False)
    axes = axes[0]
    colors_map = {"7SL": P_7SL, "U1": P_U1}
    for ax_idx, (group, length_counts) in enumerate(sorted(groups_data.items())):
        ax = axes[ax_idx]
        if not length_counts:
            ax.set_title(f"{group} group (no tail data)", fontsize=12, fontweight="bold")
            continue
        lengths = sorted(length_counts.keys())
        counts = [length_counts[l] for l in lengths]
        ax.bar(lengths, counts, 0.6, color=colors_map.get(group, P_CORAL), edgecolor="white", linewidth=0.5)
        ax.set_title(f"{group} group", fontsize=12, fontweight="bold")
        ax.set_xlabel("Tail length (nt)", fontsize=10)
        ax.set_ylabel("Read count", fontsize=10)
        ax.grid(axis="y", color="#E0E0E0", linewidth=0.5)
        ax.set_axisbelow(True)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
    fig.suptitle("Non-template Tail Length Distribution", fontsize=13, fontweight="bold", y=1.01)
    fig.tight_layout()
    return img_store.save_fig(fig, "tail_length")


def plot_end_position(samples_data: list[dict], img_store: TempImageStore) -> str:
    """Generate chart tracking the 3' end positional deviation (upstream, at_end, downstream)."""
    names = [_short_name(d["sample"]) for d in samples_data]
    groups = [infer_group(d["sample"]) for d in samples_data]
    fig, ax = plt.subplots(figsize=(8.5, 4.0))
    x = np.arange(len(names))
    w = 0.5
    upstream = _clean_list([d.get("upstream") for d in samples_data])
    at_end = _clean_list([d.get("at_end") for d in samples_data])
    downstream = _clean_list([d.get("downstream") for d in samples_data])
    ax.bar(x, upstream, w, label="Upstream (<0)", color=P_7SL)
    ax.bar(x, at_end, w, bottom=upstream, label="At end (0)", color="#5EEAD4")
    ax.bar(x, downstream, w, bottom=[u+a for u, a in zip(upstream, at_end)], label="Downstream (>0)", color=P_CORAL)
    ax.set_ylabel("Read count", fontsize=11)
    ax.set_xticks(x)
    ax.set_xticklabels(names, fontsize=8, rotation=30, ha="right")
    ax.legend(fontsize=8, framealpha=0.9)
    ax.grid(axis="y", color="#E0E0E0", linewidth=0.5)
    ax.set_axisbelow(True)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    for i in range(1, len(groups)):
        if groups[i] != groups[i-1]:
            ax.axvline(x=i - 0.5, color="#999999", linewidth=1, linestyle="--", alpha=0.5)
    fig.tight_layout()
    return img_store.save_fig(fig, "end_position")


# ══════════════════════════════════════════════════════════════════════════
# Slide builders
# ══════════════════════════════════════════════════════════════════════════

def build_title_slide(prs: Presentation, title: str, subtitle: str, date: str, pipeline_text: str, lang: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_NAVY
    _textbox(slide, 0.8, 0.85, 8.4, 0.9, title, font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        _textbox(slide, 1.1, 1.75, 7.8, 0.7, subtitle, font_size=16,
                 color=RGBColor(0xD9, 0xE7, 0xF5), align=PP_ALIGN.CENTER)
    meta_lines = []
    if date:
        meta_lines.append(f"{t('date', lang)}: {date}")
    if pipeline_text:
        meta_lines.append(f"{t('pipeline', lang)}: {pipeline_text}")
    _textbox(slide, 1.0, 2.55, 8.0, 1.25, "\n".join(meta_lines), font_size=12,
             color=RGBColor(0xC8, 0xD6, 0xE5), align=PP_ALIGN.CENTER)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(4.35), Inches(8.3), Inches(0.32))
    band.fill.solid()
    band.fill.fore_color.rgb = C_ACCENT
    band.line.fill.background()


def build_workflow_slide(prs: Presentation, pipeline_text: str, lang: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG
    _header(slide, t("workflow", lang))
    steps = [
        "FASTQ / Meta",
        "Demux + Dedup",
        "TrimGalore",
        "STAR 3-pass",
        "Gene-specific\n+ Tailer",
    ]
    box_w = 1.55
    gap = 0.2
    start_x = 0.5
    y = 1.5
    colors = [RGBColor(0xE8, 0xF1, 0xFB), RGBColor(0xE9, 0xF7, 0xF1), RGBColor(0xF9, 0xEE, 0xD7),
              RGBColor(0xF8, 0xE7, 0xEA), RGBColor(0xE9, 0xE8, 0xFA)]
    for idx, step in enumerate(steps):
        x = start_x + idx * (box_w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(0.95))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[idx]
        shape.line.color.rgb = C_ACCENT
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_NAVY
        p.alignment = PP_ALIGN.CENTER
        if idx < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + box_w + 0.03), Inches(y + 0.28), Inches(0.14), Inches(0.34))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = C_ACCENT
            arrow.line.fill.background()
    _textbox(slide, 0.6, 3.0, 8.8, 1.0, pipeline_text, font_size=13, color=C_TEXT)
    bullets = [
        ("去重通过 jla-demultiplexer 去除 PCR 重复，TrimGalore 修剪接头和质量。"
         if lang == "zh" else
         "Dedup removes PCR duplicates via jla-demultiplexer; TrimGalore trims adapters and quality."),
        ("STAR 3-pass 进行全基因组比对，提取小 RNA 相关 reads。"
         if lang == "zh" else
         "STAR 3-pass performs whole-genome alignment and extracts small-RNA-associated reads."),
        ("基因特异性比对进行 per-gene 局部重比对，Tailer 分析 3' 端非模板加尾。"
         if lang == "zh" else
         "Gene-specific alignment performs per-gene local realignment; Tailer analyses 3' non-template tails."),
    ]
    _bullets(slide, 0.7, 3.7, 8.6, 1.35, bullets, font_size=11)


def build_read_flow_slide(prs: Presentation, samples_data: list[dict], img_store: TempImageStore, lang: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG
    _header(slide, t("read_flow", lang))
    img = plot_read_flow(samples_data, img_store)
    _add_picture(slide, img, 0.3, 0.9, 6.3, 3.6)
    
    _textbox(slide, 6.8, 0.9, 3.0, 0.3, t("notes", lang), font_size=13, bold=True, color=C_NAVY)
    notes = []
    
    raw = _filter_none([d.get("raw_reads") for d in samples_data])
    dedup = _filter_none([d.get("dedup_reads") for d in samples_data])
    trim = _filter_none([d.get("trimmed_reads") for d in samples_data])
    fbam = _filter_none([d.get("final_bam_reads") for d in samples_data])
    
    if raw:
        notes.append(f"Raw: {min(raw):,}–{max(raw):,}")
    if dedup:
        notes.append(f"Post-dedup: {min(dedup):,}–{max(dedup):,}")
    if trim:
        notes.append(f"Post-trim: {min(trim):,}–{max(trim):,}")
    if fbam:
        notes.append(f"STAR BAM: {min(fbam):,}–{max(fbam):,}")
    _bullets(slide, 6.9, 1.2, 2.9, 3.3, notes, font_size=10)


def build_trimming_slide(prs: Presentation, samples_data: list[dict], img_store: TempImageStore, lang: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG
    _header(slide, t("trimming", lang))
    img = plot_trimming(samples_data, img_store)
    _add_picture(slide, img, 0.3, 0.9, 6.3, 3.6)
    _textbox(slide, 6.8, 0.9, 3.0, 0.3, t("notes", lang), font_size=13, bold=True, color=C_NAVY)
    
    r1 = _filter_none([d.get("r1_adapter_pct") for d in samples_data])
    r2 = _filter_none([d.get("r2_adapter_pct") for d in samples_data])
    pr = _filter_none([d.get("pairs_removed_pct") for d in samples_data])
    
    notes = []
    if r1: notes.append(f"R1 adapter: {min(r1):.1f}–{max(r1):.1f}%")
    if r2: notes.append(f"R2 adapter: {min(r2):.1f}–{max(r2):.1f}%")
    if pr: notes.append(f"Pairs removed: {min(pr):.1f}–{max(pr):.1f}%")
    
    notes.append("R2 adapter content higher than R1" if lang == "zh" else "R2 adapter content higher than R1")
    _bullets(slide, 6.9, 1.2, 2.9, 3.3, notes, font_size=10)


def build_star_slide(prs: Presentation, samples_data: list[dict], img_store: TempImageStore, lang: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG
    _header(slide, t("star_mapping", lang))
    img = plot_star_mapping(samples_data, img_store)
    _add_picture(slide, img, 0.3, 0.9, 9.1, 3.8)


def build_gene_slide(prs: Presentation, samples_data: list[dict], analysis_dir: str, img_store: TempImageStore, lang: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG
    _header(slide, t("gene_assignment", lang))
    img = plot_gene_assignments(samples_data, analysis_dir, img_store)
    _add_picture(slide, img, 0.3, 0.82, 9.1, 2.8)
    # Data table
    rows = [[t("sample", lang), t("genes", lang), t("total_assigned", lang), t("top_gene", lang), t("top_reads", lang)]]
    for d in samples_data:
        rows.append([_short_name(d["sample"]), d.get("n_genes"), d.get("total_assigned"), d.get("top_gene"), d.get("top_gene_reads")])
    _table(slide, 1.5, 3.8, 7.0, 0.2 * len(rows), rows, font_size=9)


def build_tail_slide(prs: Presentation, samples_data: list[dict], img_store: TempImageStore, lang: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG
    _header(slide, t("tail_analysis", lang))
    img1 = plot_tail_summary(samples_data, img_store)
    img2 = plot_end_position(samples_data, img_store)
    _add_picture(slide, img1, 0.2, 0.9, 4.8, 3.5)
    _add_picture(slide, img2, 5.1, 0.9, 4.8, 3.5)


def build_tail_length_slide(prs: Presentation, samples_data: list[dict], img_store: TempImageStore, lang: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_BG
    _header(slide, t("tail_length", lang))
    img = plot_tail_length(samples_data, img_store)
    _add_picture(slide, img, 0.5, 0.9, 9.0, 3.5)
    # Summary
    groups_data: dict[str, dict] = defaultdict(lambda: {"total": 0, "with_tail": 0})
    for d in samples_data:
        g = infer_group(d["sample"])
        groups_data[g]["total"] += (d.get("total_tail_reads") or 0)
        groups_data[g]["with_tail"] += (d.get("with_tail") or 0)
    notes = []
    for g in sorted(groups_data.keys()):
        s = groups_data[g]
        pct = s["with_tail"] / s["total"] * 100 if s["total"] > 0 else 0
        if lang == "zh":
            notes.append(f"{g} 组: {s['with_tail']} reads 有 tail ({pct:.1f}%)")
        else:
            notes.append(f"{g} group: {s['with_tail']} reads with tail ({pct:.1f}%)")
    _bullets(slide, 0.7, 4.5, 8.6, 0.8, notes, font_size=11)


def build_conclusion_slide(prs: Presentation, samples_data: list[dict], lang: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_NAVY
    _textbox(slide, 0.8, 0.75, 4.0, 0.5, t("conclusion", lang), font_size=28, bold=True, color=C_WHITE)
    bullets = []
    
    raw_total = sum(_filter_none([d.get("raw_reads") for d in samples_data]))
    dedup_total = sum(_filter_none([d.get("dedup_reads") for d in samples_data]))
    total_assigned = sum(_filter_none([d.get("total_assigned") for d in samples_data]))
    total_with_tail = sum(_filter_none([d.get("with_tail") for d in samples_data]))
    total_tail = sum(_filter_none([d.get("total_tail_reads") for d in samples_data]))

    if samples_data:
        n = len(samples_data)
        bullets.append(
            f"共 {n} 个样本完成全流程分析。" if lang == "zh" else f"{n} samples completed the full pipeline."
        )
        if raw_total > 0:
            reduction = (1 - dedup_total / raw_total) * 100
            bullets.append(
                f"去重后 reads 减少 {reduction:.1f}%。" if lang == "zh" else f"Dedup reduced reads by {reduction:.1f}%."
            )
        bullets.append(
            f"基因特异性比对共分配 {total_assigned:,} 条 reads。" if lang == "zh" else f"Gene-specific alignment assigned {total_assigned:,} reads total."
        )
        if total_tail > 0:
            pct = total_with_tail / total_tail * 100
            bullets.append(
                f"Tailer 检出 {total_with_tail:,} 条 reads 带非模板 tail ({pct:.1f}%)。" if lang == "zh"
                else f"Tailer detected {total_with_tail:,} reads with non-template tails ({pct:.1f}%)."
            )
    if not bullets:
        bullets.append("报告已生成。" if lang == "zh" else "Report generated.")
    _bullets(slide, 0.9, 1.45, 8.2, 2.35, bullets, font_size=15, color=C_WHITE)
    # Stat cards
    cards = [
        (t("sample", lang), str(len(samples_data))),
        (t("total_assigned", lang), str(total_assigned)),
        (t("with_tail", lang), str(total_with_tail)),
    ]
    x0 = 0.8
    for i, (label, value) in enumerate(cards):
        x = x0 + i * 2.75
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.0), Inches(2.45), Inches(0.7))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x25, 0x38, 0x5A)
        card.line.color.rgb = C_ACCENT
        tf = card.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = label
        p1.font.size = Pt(11)
        p1.font.color.rgb = RGBColor(0xC9, 0xD7, 0xE8)
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(20)
        p2.font.bold = True
        p2.font.color.rgb = C_WHITE
        p2.alignment = PP_ALIGN.CENTER
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.8), Inches(8.4), Inches(0.15))
    band.fill.solid()
    band.fill.fore_color.rgb = C_ACCENT
    band.line.fill.background()


# ══════════════════════════════════════════════════════════════════════════
# Excel file inventory
# ══════════════════════════════════════════════════════════════════════════

def write_file_inventory(output_path: str, analysis_dir: str, samples: List[str],
                         paired_samples: List[str], single_samples: List[str],
                         samples_data: List[dict]) -> None:
    """Generate Excel inventory containing run parameters, sample QC metrics and optional manifest details.
    Missing/None values are output as empty strings, while percentages are rendered with openpyxl '0.00%' format.
    """
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Overview"
    overview_rows = [
        ["Category", "Item", "Value"],
        ["Analysis", "analysis_dir", analysis_dir],
        ["Analysis", "n_samples", str(len(samples))],
        ["Analysis", "n_paired", str(len(paired_samples))],
        ["Analysis", "n_single", str(len(single_samples))],
        ["Analysis", "samples", ", ".join(samples)],
    ]
    for row in overview_rows:
        ws.append(row)
    ws.column_dimensions["A"].width = 12
    ws.column_dimensions["B"].width = 18
    ws.column_dimensions["C"].width = 60

    # Sample stats sheet
    ws_stats = wb.create_sheet("Sample_Stats")
    stat_cols = [
            "sample", "raw_reads", "dedup_reads", "trimmed_reads",
            "r1_adapter_counts", "r1_adapter_pct",     
            "r2_adapter_counts", "r2_adapter_pct",     
            "cutoff_threshold",  
            "pairs_removed_counts", "pairs_removed_pct", 
            "star_input_reads", "unique", "multi", "unmapped",
            "unique_pct", "multi_pct", "unmapped_pct",
            "final_bam_reads", "n_genes", "total_assigned", "top_gene",
            "top_gene_reads", "total_tail_reads",
            "with_tail", "without_tail", "upstream", "at_end", "downstream"
        ]
    
    # Write Header
    ws_stats.append(stat_cols)
    
    # Write Data Rows
    for r_idx, d in enumerate(samples_data, start=2): # Header is row 1
        for c_idx, c in enumerate(stat_cols, start=1):
            val = d.get(c)
            cell = ws_stats.cell(row=r_idx, column=c_idx)
            
            if val is None:
                cell.value = ""
            elif c.endswith("_pct") and isinstance(val, (int, float)):
                # Store pure float divided by 100 and apply openpyxl percentage formatting
                cell.value = val / 100.0
                cell.number_format = '0.00%'
            else:
                cell.value = val

    for i, col in enumerate(stat_cols, 1):
        ws_stats.column_dimensions[openpyxl.utils.get_column_letter(i)].width = min(max(len(col) + 2, 12), 30)

    # Per-sample gene manifests
    for d in samples_data:
        manifest_path = os.path.join(analysis_dir, "common", "4_per_gene_bam", d["sample"], "genes.tsv")
        df = pd.read_csv(manifest_path, sep="\t") if os.path.isfile(manifest_path) else pd.DataFrame()
        if df.empty:
            continue
        sheet_name = d["sample"][:28]
        ws_gene = wb.create_sheet(sheet_name)
        ws_gene.append([d["sample"]])
        ws_gene.append(list(df.columns))
        for _, row in df.iterrows():
            ws_gene.append([str(row[col]) if pd.notna(row[col]) else "" for col in df.columns])

    # Per-sample tail CSVs
    for d in samples_data:
        tail_path = os.path.join(analysis_dir, "common", "4_per_gene_bam", d["sample"], f"{d['sample']}_tail.csv")
        df = pd.read_csv(tail_path, skipinitialspace=True) if os.path.isfile(tail_path) else pd.DataFrame()
        if df.empty:
            continue
        if len(df) > 5000:
            df = df.head(5000)
        sheet_name = f"tail_{d['sample']}"[:28]
        ws_tail = wb.create_sheet(sheet_name)
        ws_tail.append([d["sample"]])
        ws_tail.append(list(df.columns))
        for _, row in df.iterrows():
            ws_tail.append([str(row[col]) if pd.notna(row[col]) else "" for col in df.columns])

    wb.save(output_path)


# ══════════════════════════════════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════════════════════════════════

def main():
    ap = argparse.ArgumentParser(description="Generate ncRNAseq PPT report")
    ap.add_argument("--analysis-dir", required=True)
    ap.add_argument("--output", required=True)
    ap.add_argument("--samples", nargs="*", default=[])
    ap.add_argument("--paired-samples", nargs="*", default=[])
    ap.add_argument("--single-samples", nargs="*", default=[])
    ap.add_argument("--title", default="")
    ap.add_argument("--subtitle", default="")
    ap.add_argument("--pipeline", default="")
    ap.add_argument("--genome", default="")
    ap.add_argument("--date", default="")
    ap.add_argument("--lang", default="zh")
    ap.add_argument("--img-dir", default="")
    ap.add_argument("--file-inventory", default="")
    args = ap.parse_args()

    analysis_dir = os.path.abspath(args.analysis_dir)
    output = os.path.abspath(args.output)
    os.makedirs(os.path.dirname(output), exist_ok=True)

    samples = args.samples
    if not samples:
        # Auto-detect from 4_per_gene_bam
        pg_dir = os.path.join(analysis_dir, "common", "4_per_gene_bam")
        if os.path.isdir(pg_dir):
            samples = sorted([d for d in os.listdir(pg_dir)
                              if os.path.isdir(os.path.join(pg_dir, d))])
    if not samples:
        # Fallback: 3_raw_bam
        bam_dir = os.path.join(analysis_dir, "common", "3_raw_bam")
        if os.path.isdir(bam_dir):
            samples = sorted([d for d in os.listdir(bam_dir)
                              if os.path.isdir(os.path.join(bam_dir, d))])
    print(f"Detected samples: {samples}", file=sys.stderr)

    samples_data = collect_sample_data(analysis_dir, samples)
    if not samples_data:
        print("WARNING: No sample data collected", file=sys.stderr)

    pipeline_text = args.pipeline or "FASTQ -> Demux/Dedup -> TrimGalore -> STAR 3-pass -> Gene-specific alignment + Tailer"
    title = args.title or t("title", args.lang)
    subtitle_parts = []
    if args.subtitle:
        subtitle_parts.append(args.subtitle)
    if args.genome:
        subtitle_parts.append(args.genome)
    subtitle = " | ".join(subtitle_parts)

    img_store = TempImageStore(args.img_dir)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    build_title_slide(prs, title, subtitle, args.date, pipeline_text, args.lang)
    build_workflow_slide(prs, pipeline_text, args.lang)

    if samples_data:
        build_read_flow_slide(prs, samples_data, img_store, args.lang)
        build_trimming_slide(prs, samples_data, img_store, args.lang)
        build_star_slide(prs, samples_data, img_store, args.lang)
        build_gene_slide(prs, samples_data, analysis_dir, img_store, args.lang)
        build_tail_slide(prs, samples_data, img_store, args.lang)
        build_tail_length_slide(prs, samples_data, img_store, args.lang)
        build_conclusion_slide(prs, samples_data, args.lang)

    prs.save(output)
    img_store.cleanup()
    print(f"Report saved to {output}", file=sys.stderr)
    if args.file_inventory:
        inventory_path = os.path.abspath(args.file_inventory)
        os.makedirs(os.path.dirname(inventory_path), exist_ok=True)
        write_file_inventory(inventory_path, analysis_dir, samples,
                             args.paired_samples, args.single_samples, samples_data)
        print(f"data is saved to {inventory_path}", file=sys.stderr)

if __name__ == "__main__":
    main()
